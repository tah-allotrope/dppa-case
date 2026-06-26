from pathlib import Path
import shutil

from pptx import Presentation


DECK = Path("ceba/DPPA Presentation July 2026 Scenario Training.pptx")
BACKUP = Path("ceba/DPPA Presentation July 2026 Scenario Training.backup-2026-06-26.pptx")
VERIFY = Path("deck-qa/july-deck-corrections-verify.txt")


REPLACEMENTS = {
    5: [
        ("8,263,196,000", "8,563,196,000"),
    ],
    6: [
        ("8,263,196,000 + 500,000,000 = 8,763,196,000", "8,563,196,000 + 500,000,000 = 9,063,196,000"),
        ("8,763,196,000", "9,063,196,000"),
    ],
    7: [
        ("8,2 billion", "8.5 billion"),
        ("8.2 billion", "8.5 billion"),
        ("6.2 billion", "5.8 billion"),
    ],
}

TABLE_REPLACEMENTS = {
    "1,026": "1.026",
    "1,008": "1.008",
    "163,3": "163.30",
}


def replace_in_paragraph(paragraph, old, new):
    text = "".join(run.text for run in paragraph.runs)
    if old not in text:
        return 0

    replaced = text.replace(old, new, 1)
    remaining = replaced
    for run in paragraph.runs:
        run_len = len(run.text)
        run.text = remaining[:run_len]
        remaining = remaining[run_len:]
    if remaining and paragraph.runs:
        paragraph.runs[-1].text += remaining
    return 1


def iter_paragraphs(shape):
    if getattr(shape, "has_text_frame", False):
        yield from shape.text_frame.paragraphs
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            for cell in row.cells:
                for paragraph in cell.text_frame.paragraphs:
                    yield paragraph


def apply_slide_replacements(prs):
    counts = {}
    for slide_no, pairs in REPLACEMENTS.items():
        slide = prs.slides[slide_no - 1]
        for old, new in pairs:
            count = 0
            for shape in slide.shapes:
                for paragraph in iter_paragraphs(shape):
                    count += replace_in_paragraph(paragraph, old, new)
            counts[(slide_no, old)] = count
    return counts


def apply_slide_3_table_replacements(prs):
    counts = {old: 0 for old in TABLE_REPLACEMENTS}
    slide = prs.slides[2]
    for shape in slide.shapes:
        if not getattr(shape, "has_table", False):
            continue
        for row in shape.table.rows:
            for cell in row.cells:
                for paragraph in cell.text_frame.paragraphs:
                    for old, new in TABLE_REPLACEMENTS.items():
                        counts[old] += replace_in_paragraph(paragraph, old, new)
    return counts


def dump_text(prs):
    lines = []
    for idx, slide in enumerate(prs.slides, start=1):
        lines.append(f"--- Slide {idx} ---")
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = shape.text.strip()
                if text:
                    lines.append(text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    if any(cells):
                        lines.append(" | ".join(cells))
    return "\n".join(lines) + "\n"


def main():
    if not DECK.exists():
        raise FileNotFoundError(DECK)
    if not BACKUP.exists():
        shutil.copy2(DECK, BACKUP)

    prs = Presentation(DECK)
    if len(prs.slides) != 11:
        raise RuntimeError(f"Expected 11 slides, found {len(prs.slides)}")

    counts = apply_slide_replacements(prs)
    table_counts = apply_slide_3_table_replacements(prs)

    missing = [f"slide {slide_no}: {old}" for (slide_no, old), count in counts.items() if count == 0]
    missing.extend([f"slide 3 table: {old}" for old, count in table_counts.items() if count == 0])
    if missing:
        raise RuntimeError("Missing expected replacements:\n" + "\n".join(missing))

    prs.save(DECK)
    reloaded = Presentation(DECK)
    if len(reloaded.slides) != 11:
        raise RuntimeError(f"Reloaded deck slide count changed: {len(reloaded.slides)}")

    VERIFY.parent.mkdir(parents=True, exist_ok=True)
    VERIFY.write_text(dump_text(reloaded), encoding="utf-8")

    print("Deck corrections applied")
    print(f"Backup: {BACKUP}")
    print(f"Verify: {VERIFY}")
    print("Slides:", len(reloaded.slides))


if __name__ == "__main__":
    main()
