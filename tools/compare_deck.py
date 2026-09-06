"""
compare_deck.py
================

Guard (PHASE-05 of plans/2026-08-22-delivery-stall-recovery-plan.md): compares
two .pptx files at the text level and reports the first difference, so a
freshly-built deck can be checked against the committed one without a human
eyeballing 27 slides.

Binary .pptx files cannot be diffed directly -- each save re-serializes zip
member ordering, timestamps, and relationship IDs even when the visible
content is unchanged, so a byte-for-byte comparison would report differences
on every rebuild regardless of whether anything a viewer would notice
actually changed. Text-level comparison is the workable form: it extracts
every body shape's text and every slide's speaker notes, in document order,
and reports the first slide/origin where the two decks disagree.

Run:  python tools/compare_deck.py A.pptx B.pptx
      (Windows, if the default python is shadowed: PYTHONPATH= py tools/compare_deck.py A.pptx B.pptx)
Exit 0 + "IDENTICAL" if every slide's body and notes text matches; exit 1 +
the first differing (slide, origin, text_a, text_b) otherwise.
"""
from __future__ import annotations

import sys

from pptx import Presentation


def extract_text(pptx_path: str) -> list[tuple[int, str, str]]:
    """Return (slide_index, origin, text) triples for every body shape's text
    frame and, when present, each slide's speaker notes -- in document order.
    origin is "body" or "notes"."""
    prs = Presentation(pptx_path)
    triples: list[tuple[int, str, str]] = []
    for slide_index, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_text_frame:
                triples.append((slide_index, "body", shape.text_frame.text))
        if slide.has_notes_slide:
            triples.append((slide_index, "notes", slide.notes_slide.notes_text_frame.text))
    return triples


def first_difference(
    a: list[tuple[int, str, str]], b: list[tuple[int, str, str]]
) -> tuple[int, str, str, str] | None:
    """Return (slide_index, origin, text_a, text_b) for the first entry where a
    and b disagree -- either a differing text at the same position, or one
    list running out before the other. Return None when the two lists are
    identical."""
    for i in range(max(len(a), len(b))):
        entry_a = a[i] if i < len(a) else None
        entry_b = b[i] if i < len(b) else None
        if entry_a is None:
            return (entry_b[0], entry_b[1], "<no shape>", entry_b[2])
        if entry_b is None:
            return (entry_a[0], entry_a[1], entry_a[2], "<no shape>")
        if entry_a != entry_b:
            slide_index, origin, text_a = entry_a
            _, _, text_b = entry_b
            return (slide_index, origin, text_a, text_b)
    return None


def main(argv: list[str] | None = None) -> int:
    if hasattr(sys.stdout, "reconfigure"):
        sys.stdout.reconfigure(encoding="utf-8", errors="replace")
    argv = argv if argv is not None else sys.argv[1:]
    if len(argv) != 2:
        print("Usage: python tools/compare_deck.py A.pptx B.pptx", file=sys.stderr)
        return 2

    path_a, path_b = argv
    text_a = extract_text(path_a)
    text_b = extract_text(path_b)

    diff = first_difference(text_a, text_b)
    if diff is None:
        print(f"COMPARE-DECK: IDENTICAL ({len(text_a)} body/notes entries across both decks)")
        return 0

    slide_index, origin, value_a, value_b = diff
    print(f"COMPARE-DECK: DIFFERS at slide {slide_index} [{origin}]")
    print(f"  {path_a}: {value_a!r}")
    print(f"  {path_b}: {value_b!r}")
    return 1


if __name__ == "__main__":
    sys.exit(main(sys.argv[1:]))
