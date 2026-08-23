"""Unit tests for tools/compare_deck.py (PHASE-05 of
plans/2026-08-22-delivery-stall-recovery-plan.md)."""
from __future__ import annotations

import sys
import tempfile
import unittest
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

TOOLS_DIR = Path(__file__).resolve().parent.parent
if str(TOOLS_DIR) not in sys.path:
    sys.path.insert(0, str(TOOLS_DIR))

import compare_deck as cd  # noqa: E402


def _build_deck(path: Path, slide_texts: list[tuple[str, str | None]]) -> None:
    """Build a minimal .pptx at `path`. Each entry is (body_text, notes_text_or_None)."""
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]
    for body_text, notes_text in slide_texts:
        slide = prs.slides.add_slide(blank_layout)
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        box.text_frame.text = body_text
        if notes_text is not None:
            slide.notes_slide.notes_text_frame.text = notes_text
    prs.save(str(path))


class TestExtractText(unittest.TestCase):
    def setUp(self):
        self.tmp = tempfile.TemporaryDirectory()
        self.addCleanup(self.tmp.cleanup)

    def test_extracts_body_and_notes_in_order(self):
        path = Path(self.tmp.name) / "deck.pptx"
        _build_deck(path, [("Slide one body", "Slide one notes"), ("Slide two body", None)])
        triples = cd.extract_text(str(path))
        self.assertEqual(
            triples,
            [
                (0, "body", "Slide one body"),
                (0, "notes", "Slide one notes"),
                (1, "body", "Slide two body"),
            ],
        )


class TestFirstDifference(unittest.TestCase):
    def test_identical_lists_return_none(self):
        a = [(0, "body", "x"), (1, "notes", "y")]
        self.assertIsNone(cd.first_difference(a, list(a)))

    def test_differing_body_text_is_reported(self):
        a = [(0, "body", "old text")]
        b = [(0, "body", "new text")]
        self.assertEqual(cd.first_difference(a, b), (0, "body", "old text", "new text"))

    def test_shorter_list_is_reported_as_a_difference(self):
        a = [(0, "body", "x"), (1, "body", "y")]
        b = [(0, "body", "x")]
        result = cd.first_difference(a, b)
        self.assertEqual(result[0], 1)
        self.assertEqual(result[1], "body")
        self.assertEqual(result[3], "<no shape>")


class TestMain(unittest.TestCase):
    def setUp(self):
        self.tmp = tempfile.TemporaryDirectory()
        self.addCleanup(self.tmp.cleanup)

    def test_identical_decks_exit_zero(self):
        path_a = Path(self.tmp.name) / "a.pptx"
        path_b = Path(self.tmp.name) / "b.pptx"
        _build_deck(path_a, [("Same slide text", None)])
        _build_deck(path_b, [("Same slide text", None)])
        self.assertEqual(cd.main([str(path_a), str(path_b)]), 0)

    def test_planted_change_is_detected(self):
        """A build whose only change is one run of text on one slide is caught."""
        path_a = Path(self.tmp.name) / "a.pptx"
        path_b = Path(self.tmp.name) / "b.pptx"
        _build_deck(
            path_a,
            [("Slide 1", "notes 1"), ("Slide 2 unchanged", "notes 2"), ("Slide 3", None)],
        )
        _build_deck(
            path_b,
            [("Slide 1", "notes 1"), ("Slide 2 CHANGED", "notes 2"), ("Slide 3", None)],
        )
        self.assertEqual(cd.main([str(path_a), str(path_b)]), 1)

    def test_wrong_argument_count_exits_two(self):
        self.assertEqual(cd.main(["only-one-arg.pptx"]), 2)


if __name__ == "__main__":
    unittest.main()
