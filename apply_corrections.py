"""
Post-processing corrections for the consolidated deck.

P0-3: Remove slide 23 (old Cases 5/6 intro bridge) — has off-slide table
      and references removed Case Studies 5/6.

P1:   Move the 5 new slides (3 canonical cases + 2 callouts) under the
      Module 5 divider (slide 22) so the audience sees the cases right
      after "Three Canonical Cases" instead of after Module 6 wrap.

Order after fix:
  22: Module 5 divider (unchanged)
  23: Canonical Case A (from slide 35)
  24: Canonical Case B (from slide 36)
  25: Canonical Case C (from slide 37)
  26: Financing callout (from slide 24)
  27: Netting callout (from slide 38)
  28+: Module 6 wrap, Q&A, panel, bios, interactive, formulas, recap, wrap
"""
from pptx import Presentation

SRC = "ceba/CEBA DPPA 2026.pptx"
OUT = "ceba/CEBA DPPA 2026.pptx"
SLIDE_BRIDGE = 22  # 0-based = 1-based slide 23


def _find_sldId_for_slide(prs, slide):
    sldIdLst = prs.slides._sldIdLst
    for child in sldIdLst:
        rid = child.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
        try:
            rel = prs.part.rels[rid]
        except KeyError:
            continue
        target = getattr(rel, "target_part", None)
        if target is slide.part:
            return child
    return None


def reorder_slide(prs, slide, target_idx_0based):
    sldIdLst = prs.slides._sldIdLst
    sldId_el = _find_sldId_for_slide(prs, slide)
    if sldId_el is None:
        raise RuntimeError("Could not find sldId for slide")
    sldIdLst.remove(sldId_el)
    children = list(sldIdLst)
    if target_idx_0based >= len(children):
        sldIdLst.append(sldId_el)
    else:
        sldIdLst.insert(children.index(children[target_idx_0based]), sldId_el)


def delete_slide_by_index(prs, idx_0based):
    sldIdLst = prs.slides._sldIdLst
    children = list(sldIdLst)
    if idx_0based >= len(children):
        raise IndexError(f"Index {idx_0based} out of range")
    target_sldId = children[idx_0based]
    rid = target_sldId.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
    prs.part.drop_rel(rid)
    sldIdLst.remove(target_sldId)


def find_slide_0based(prs, title_substring):
    """Return the 0-based index of the first slide whose text contains the
       given substring, or -1 if not found."""
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if hasattr(shape, "text") and title_substring in shape.text:
                return i
    return -1


def main():
    prs = Presentation(SRC)
    print(f"Opening deck: {len(prs.slides)} slides\n")

    # --- P0-3: Remove slide 23 (old Cases 5/6 bridge) ---
    delete_slide_by_index(prs, SLIDE_BRIDGE)
    print(f"After removing slide 23 (bridge): {len(prs.slides)} slides\n")

    # --- P1: Reorder the 5 new slides under Module 5 divider ---
    # Each reorder_slide call shifts other elements; we must
    # re-find the slide by title after each move because its
    # 0-based index has changed.
    to_reorder = [
        ("CANONICAL CASE A", 22),   # → 1-based 23
        ("CANONICAL CASE B", 23),   # → 1-based 24
        ("CANONICAL CASE C", 24),   # → 1-based 25
        ("DEVELOPER FINANCING", 25),  # → 1-based 26
        ("MULTI-PARTY NETTING", 26),   # → 1-based 27
    ]
    for substring, target in to_reorder:
        idx = find_slide_0based(prs, substring)
        if idx < 0:
            print(f"  WARNING: slide containing {substring!r} not found!")
            continue
        slide = prs.slides[idx]
        reorder_slide(prs, slide, target)
        print(f"  Moved '{substring}' from 1-based {idx+1} to 1-based {target+1}")

    print(f"\nFinal deck: {len(prs.slides)} slides")
    print("Layout 21-33:")
    for i in range(20, min(33, len(prs.slides))):
        slide = prs.slides[i]
        title = ""
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                title = shape.text.strip().split("\n")[0][:70]
                break
        print(f"  {i+1:2d}. {title}")

    prs.save(OUT)
    print(f"\nSaved to {OUT}")


if __name__ == "__main__":
    main()
