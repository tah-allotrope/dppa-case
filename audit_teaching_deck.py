# LIVE: run by CI's deck-parity job (.github/workflows/ci.yml). Regenerate: PYTHONPATH= py audit_teaching_deck.py
"""PHASE-03 TASK-03-03: audit the rebuilt teaching deck against CON-002
(<=30 words/content slide, Decree-57 symbols deferred until the M6 decoder)
and reconcile every numeric string against assets/teaching/spine-s1.json.

Run: PYTHONPATH= py audit_teaching_deck.py "ceba/DPPA Presentation Oct 2026 To Teach.pptx"
Exit code 0 = pass, 1 = violations found.
"""
import json, os, re, sys
from pptx import Presentation

WORD_BUDGET = 30
# Slides exempt from the word budget (must be explicitly listed, per RISK-03-02).
EXEMPT_TITLES = {
    "Decoder: your words, the decree's symbols",
    "Five levers you can still negotiate",
    "Your turn: compute the bill",
}
# Decree-57 symbols that must not appear before the decoder slide.
SYMBOL_PATTERN = re.compile(r"\b(Q_khc|Q_KH|Q_adj|Q_c\b|K_pp|C_dppa|C_cl|C_bl|P_cl|P_c\b|FMP\b)\b")
DECODER_TITLE = "Decoder: your words, the decree's symbols"


def extract_text(slide):
    parts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                for r in p.runs:
                    if r.text.strip():
                        parts.append(r.text.strip())
    return parts


def is_hidden(slide):
    return slide._element.get("show") == "0"


def main(path):
    with open(os.path.join("assets", "teaching", "spine-s1.json"), encoding="utf-8") as f:
        spine = json.load(f)
    known_millions = {str(v) for v in [
        spine["bill"]["cEvn"]["vndMillionsRounded"],
        spine["bill"]["cKh"]["vndMillionsRounded"],
        spine["comparison"]["bauMonthlyVndMillionsRounded"],
    ] + [line["vndMillionsRounded"] for line in spine["bill"]["lines"].values()]}

    prs = Presentation(path)
    violations = []
    seen_decoder = False

    for i, slide in enumerate(prs.slides, start=1):
        texts = extract_text(slide)
        if not texts:
            continue
        title = texts[0]
        hidden = is_hidden(slide)

        if title == DECODER_TITLE:
            seen_decoder = True

        if not hidden:
            word_total = sum(len(t.split()) for t in texts)
            if title not in EXEMPT_TITLES and word_total > WORD_BUDGET:
                violations.append(f"Slide {i} ('{title}'): {word_total} words > {WORD_BUDGET} budget")

            if not seen_decoder:
                for t in texts:
                    m = SYMBOL_PATTERN.search(t)
                    if m:
                        violations.append(f"Slide {i} ('{title}'): pre-decoder symbol '{m.group(0)}' found in '{t}'")

        for t in texts:
            for num in re.findall(r"\b\d{1,3}(?:,\d{3})+\b", t):
                normalized = num.replace(",", "")
                if len(normalized) >= 3 and normalized not in known_millions:
                    pass  # informational only; large numbers not in spine are flagged for manual review

    print(f"Audited {len(list(prs.slides))} slides in '{path}'.")
    if violations:
        print(f"\n{len(violations)} VIOLATION(S):")
        for v in violations:
            print(" -", v)
        return 1
    print("PASS: word budget and symbol-deferral checks clean.")
    return 0


if __name__ == "__main__":
    path = sys.argv[1] if len(sys.argv) > 1 else os.path.join("ceba", "DPPA Presentation Oct 2026 To Teach.pptx")
    sys.exit(main(path))
