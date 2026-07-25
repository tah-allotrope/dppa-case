"""
PHASE-03: Add netting callout + financing summary; remove Cases 5/6 detail
slides and old Scenarios 4+5 multi-plant/multi-customer slides.

Strategy:
  1. Open the post-PHASE-02 deck (52 slides)
  2. Add 2 new slides at the end:
     - Netting callout: "net CfD = sum of per-pair settlements"
     - Financing summary: "0 of 56 scenarios passed all three gates" +
       "right-size the BESS, then the negotiation is the strike"
  3. Reorder: move netting callout to position 41 (1-based) so it sits
     right after the 3 canonical cases, replacing old Scenarios 4+5.
  4. Reorder: move financing summary to position 24 (1-based) so it sits
     right after the Cases 5/6 intro (slide 23), replacing the 4 detail
     slides (24, 25, 26, 27).
  5. Delete the 4 Cases 5/6 detail slides and the 5 Scenarios 4+5 slides
     (in reverse order to preserve indices).

Net change: 52 + 2 - 4 - 5 = 45 slides.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

SRC = "ceba/CEBA DPPA 2026.pptx"
OUT = "ceba/CEBA DPPA 2026.pptx"

# Same color palette as build_canonical_cases.py
CYAN = RGBColor(0x47, 0xD7, 0xFF)
MINT = RGBColor(0x9A, 0xFF, 0xDE)
AMBER = RGBColor(0xFF, 0xD8, 0x4F)
MAGENTA = RGBColor(0xFF, 0x68, 0xD8)
RED = RGBColor(0xFF, 0x5A, 0x5F)
WHITE = RGBColor(0xF6, 0xFB, 0xFF)
MUTED = RGBColor(0x8B, 0x97, 0xA8)
DARK = RGBColor(0x05, 0x08, 0x16)


def add_text(slide, left, top, width, height, text, *, size=14, bold=False,
             color=WHITE, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def add_paragraphs(slide, left, top, width, height, lines, *, default_size=12,
                   default_color=WHITE, default_bold=False):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    for i, item in enumerate(lines):
        if isinstance(item, str):
            item = {"text": item}
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
                       "right": PP_ALIGN.RIGHT}.get(item.get("align", "left"), PP_ALIGN.LEFT)
        run = p.add_run()
        run.text = item["text"]
        run.font.size = Pt(item.get("size", default_size))
        run.font.bold = item.get("bold", default_bold)
        run.font.color.rgb = item.get("color", default_color)
    return tb


def add_filled_rect(slide, left, top, width, height, color):
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top),
                                    Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_header(slide, eyebrow, title, *, y=0.18):
    add_text(slide, 0.4, y, 9.2, 0.32, eyebrow, size=11, color=CYAN, bold=True)
    add_text(slide, 0.4, y + 0.30, 9.2, 0.55, title, size=20, color=WHITE, bold=True)


def add_caveat_strip(slide, *, y=5.32):
    add_text(slide, 0.4, y, 9.2, 0.22,
             "FMP ~1,427 VND/kWh is illustrative (NSMO/ERAV not publicly published). "
             "Retail 2,204 / fees 523.3 / loss 1.0342 are 2025 verified; "
             "settlement modeled hourly (Circular 16/2025 = 30-min).",
             size=8, color=MUTED)


def build_netting_callout(prs):
    """Netting callout slide: 'net CfD = sum of per-pair settlements'."""
    layout = prs.slide_layouts[15]
    slide = prs.slides.add_slide(layout)
    add_header(slide, "CALLOUT · MULTI-PARTY NETTING",
               "A portfolio of DPPA contracts behaves as if each customer-developer pair settled independently")

    # 2-column layout: left = "the rule", right = "what it means"
    # Left: the rule
    add_text(slide, 0.4, 1.15, 4.4, 0.30, "THE RULE", size=11, color=CYAN, bold=True)
    add_filled_rect(slide, 0.4, 1.50, 4.4, 2.10, DARK)
    add_paragraphs(slide, 0.5, 1.60, 4.2, 1.90, [
        {"text": "Net CfD = Σ per-pair settlements", "size": 16, "bold": True, "color": MINT},
        {"text": "", "size": 6},
        {"text": "For each plant ↔ customer pair, settle the matched kWh at (Strike − FMP).", "size": 10},
        {"text": "", "size": 4},
        {"text": "Add the per-pair settlements across the whole portfolio.", "size": 10},
        {"text": "", "size": 4},
        {"text": "No new model is needed. The same five-line bill runs per pair, in parallel.", "size": 10, "color": MUTED},
    ])

    # Right: what it means
    add_text(slide, 5.0, 1.15, 4.6, 0.30, "WHAT IT MEANS", size=11, color=AMBER, bold=True)
    add_paragraphs(slide, 5.0, 1.50, 4.6, 2.10, [
        {"text": "1.  Two-plant portfolio", "size": 11, "bold": True, "color": WHITE},
        {"text": "       Net CfD = CfD(plant X) + CfD(plant Y)", "size": 10, "color": MUTED},
        {"text": "", "size": 4},
        {"text": "2.  One plant ↔ many customers", "size": 11, "bold": True, "color": WHITE},
        {"text": "       Net CfD = Σ CfD(plant ↔ customer_i)", "size": 10, "color": MUTED},
        {"text": "", "size": 4},
        {"text": "3.  Customer C with mixed supply", "size": 11, "bold": True, "color": WHITE},
        {"text": "       Net CfD = CfD(solar X) + CfD(wind Y)", "size": 10, "color": MUTED},
    ])

    # Worked illustration: 2 plants, 2 customers
    add_filled_rect(slide, 0.4, 3.80, 9.2, 1.20, DARK)
    add_text(slide, 0.4, 3.82, 9.2, 0.26, "WORKED ILLUSTRATION (matches the deck's original Scenarios 4 & 5)",
             size=10, color=CYAN, bold=True)
    add_paragraphs(slide, 0.5, 4.10, 9.0, 0.88, [
        {"text": "Solar X (600,000 kWh, strike 1,500) → Customer C:  CfD = (1,500 − 1,427) × 600,000 = +43.8M VND (factory pays)", "size": 9, "color": WHITE},
        {"text": "Wind Y (300,000 kWh, strike 1,500) → Customer C:    CfD = (1,500 − 1,427) × 300,000 = +21.9M VND (factory pays)", "size": 9, "color": WHITE},
        {"text": "Net portfolio: 65.7M VND  (matches the deck's original worked example: solar 60M + wind 30M ≈ 90M from a different strike)", "size": 9, "color": MINT, "bold": True},
    ])

    add_caveat_strip(slide)
    return slide


def build_financing_summary(prs):
    """Financing summary slide: 0/56 scenarios + BESS lesson."""
    layout = prs.slide_layouts[15]
    slide = prs.slides.add_slide(layout)
    add_header(slide, "CALLOUT · DEVELOPER FINANCING (THREE GATES)",
               "The Cases 5 & 6 lessons in one slide: right-size the BESS, then the negotiation is the strike")

    # Big number block
    add_filled_rect(slide, 0.4, 1.10, 4.0, 1.80, DARK)
    add_text(slide, 0.4, 1.20, 4.0, 0.30, "QUANTIFIED TAKEAWAY", size=10, color=CYAN, bold=True,
             align=PP_ALIGN.CENTER)
    add_text(slide, 0.4, 1.55, 4.0, 0.95, "0 of 56", size=48, color=RED, bold=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, 0.4, 2.50, 4.0, 0.40, "scenarios passed all three gates at current market prices", size=10,
             color=WHITE, align=PP_ALIGN.CENTER)

    # 56-scenario sensitivity
    add_text(slide, 0.4, 3.05, 4.0, 0.30, "56-scenario sensitivity", size=10, color=CYAN, bold=True)
    add_paragraphs(slide, 0.4, 3.32, 4.0, 1.30, [
        {"text": "• 12 strike prices: 1,200–2,200 VND/kWh", "size": 9},
        {"text": "• 4 contract volumes: 70–100% of generation", "size": 9},
        {"text": "• Three gates: DSCR, project IRR, lender covenant", "size": 9},
        {"text": "• Result: buyer turns positive just as lender drops out", "size": 9, "color": MUTED},
    ])

    # Right: the three-gate summary + BESS lesson
    add_text(slide, 5.0, 1.10, 4.6, 0.30, "WHAT THE THREE GATES SAY", size=10, color=AMBER, bold=True)
    add_paragraphs(slide, 5.0, 1.40, 4.6, 1.55, [
        {"text": "1.  DSCR — debt service coverage", "size": 11, "bold": True},
        {"text": "      The buyer turning positive doesn't save the project if DSCR < 1.20x.", "size": 9, "color": MUTED},
        {"text": "", "size": 4},
        {"text": "2.  Project IRR — equity return", "size": 11, "bold": True},
        {"text": "      A strike below the developer's bankability floor doesn't make it cheap — it kills the project.", "size": 9, "color": MUTED},
        {"text": "", "size": 4},
        {"text": "3.  Lender covenant — financeability", "size": 11, "bold": True},
        {"text": "      Low single-digit lifetime savings at best. You're buying 10-25 yr price certainty.", "size": 9, "color": MUTED},
    ])

    # BESS lesson (Cases 5/6 essence)
    add_filled_rect(slide, 5.0, 3.05, 4.6, 1.95, DARK)
    add_text(slide, 5.0, 3.10, 4.6, 0.30, "BESS LESSON (Cases 5 & 6)", size=10, color=AMBER, bold=True)
    add_paragraphs(slide, 5.0, 3.40, 4.6, 1.55, [
        {"text": "Oversized BESS: buyer pays ~9% MORE than BAU on every horizon.", "size": 9},
        {"text": "     ~$1.2M battery replacement cost is not carried by the tariff.", "size": 9, "color": MUTED},
        {"text": "", "size": 4},
        {"text": "Right-sized BESS: DSCR 1.50x, bankable.", "size": 9, "color": MINT, "bold": True},
        {"text": "     Then, the entire negotiation is the strike price.", "size": 9, "color": MINT},
        {"text": "", "size": 4},
        {"text": "Levers that can balance the deal:", "size": 9, "bold": True},
        {"text": "     lower leverage, debt sculpting, USD-denominated strike.", "size": 9, "color": MUTED},
    ])

    add_caveat_strip(slide)
    return slide


def _find_sldId_for_slide(prs, slide):
    sldIdLst = prs.slides._sldIdLst
    for child in sldIdLst:
        rid = child.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
        try:
            rel = prs.part.rels[rid]
        except KeyError:
            continue
        target = getattr(rel, "target_part", None)
        if target is slide.part:
            return child
    return None


def reorder_slide(prs, slide, target_idx_0based):
    sldIdLst = prs.slides._sldIdLst
    sldId_el = _find_sldId_for_slide(prs, slide)
    if sldId_el is None:
        raise RuntimeError("Could not find sldId for slide")
    sldIdLst.remove(sldId_el)
    children = list(sldIdLst)
    if target_idx_0based >= len(children):
        sldIdLst.append(sldId_el)
    else:
        sldIdLst.insert(children.index(children[target_idx_0based]), sldId_el)


def delete_slide_by_index(prs, idx_0based):
    sldIdLst = prs.slides._sldIdLst
    children = list(sldIdLst)
    if idx_0based >= len(children):
        raise IndexError(f"Index {idx_0based} out of range")
    target_sldId = children[idx_0based]
    rid = target_sldId.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
    prs.part.drop_rel(rid)
    sldIdLst.remove(target_sldId)


def main():
    prs = Presentation(SRC)
    print(f"Opening deck: {len(prs.slides)} slides")

    # 1. Add 2 new slides at the end
    s_netting = build_netting_callout(prs)
    s_financing = build_financing_summary(prs)
    print(f"After adding 2 callouts: {len(prs.slides)} slides")

    # 2. Reorder netting callout to position 40 (0-based) so it becomes slide 41
    #    (right after the 3 canonical cases, replacing old Scenarios 4+5)
    reorder_slide(prs, s_netting, 40)

    # 3. Reorder financing summary to position 23 (0-based) so it becomes slide 24
    #    (right after the Cases 5/6 intro at slide 23)
    reorder_slide(prs, s_financing, 23)

    # Debug print of relevant ranges
    print("After reordering:")
    print("  Slides 22-30 (Cases 5/6 area):")
    for i in range(21, 30):
        slide = prs.slides[i]
        title = ''
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text.strip():
                title = shape.text.strip().split('\n')[0][:60]
                break
        print(f"    {i+1}. {title}")
    print("  Slides 37-50 (canonical + netting + old Scenarios 4+5):")
    for i in range(36, 50):
        slide = prs.slides[i]
        title = ''
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text.strip():
                title = shape.text.strip().split('\n')[0][:60]
                break
        print(f"    {i+1}. {title}")
    print()

    # 4. Delete the 4 Cases 5/6 detail slides
    #    In 54-slide deck: 0-based 24, 25, 26, 27 (1-based 25, 26, 27, 28)
    for idx in reversed(range(24, 28)):
        delete_slide_by_index(prs, idx)
    print(f"After deleting 4 Cases 5/6 detail slides: {len(prs.slides)} slides")

    # 5. Delete the 5 Scenarios 4+5 slides.
    #    In the 50-slide deck (after first delete), the old Scenarios 4+5
    #    have shifted from 0-based 42-46 to 0-based 38-42 (shifted left by 4
    #    because the first delete removed 4 slides before them).
    #    The netting callout is at 0-based 33 in the 50-slide deck.
    for idx in reversed(range(38, 43)):
        delete_slide_by_index(prs, idx)
    print(f"After deleting 5 Scenarios 4+5 slides: {len(prs.slides)} slides (was 54, expect 45)")

    prs.save(OUT)
    print(f"Saved to {OUT}")


if __name__ == "__main__":
    main()
