from pptx import Presentation
import sys

for path in sys.argv[1:]:
    print(f"\n=== {path} ===")
    try:
        prs = Presentation(path)
        print(f"Slides: {len(prs.slides)}")
        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip()[:120])
            print(f"\nSlide {i+1}:")
            for t in texts[:8]:
                print(f"  - {t}")
    except Exception as e:
        print(f"Error: {e}")
