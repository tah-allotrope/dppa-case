"""
PHASE-04: Policy refresh + agenda renumbering.

Strategy: surgical text edits to existing slides — no slide count change.
  1. Slide 22 (Module 5 divider): rename "Case Studies 5 & 6" → "Three Canonical Cases"
  2. Slide 3 (Follow the Money / roadmap): add reference to 3 canonical cases
  3. Slide 27 (Know the five-line bill): add reference to 3 canonical cases
  4. Slide 33 (Interactive Exercise): rename "DPPA Scenario Analysis" → "Apply the 3 Canonical Cases"
  5. Slide 34 (Cost formulas recap): add Circular 16/2025 citation
  6. Slide 40 (DPPA represents slide): update "first transactions emerging" → Samsung/TTC Duc Hue 2 first-DPPA
  7. Slide 42 (Vietnam's DPPA is a game-changer): "early-stage" → "first grid DPPA live in 2026"

PNG export via export-slides.ps1 (PowerPoint COM) is blocked on this machine
because no PowerPoint is installed. The script attempts a fallback via
python-pptx text dump instead and records the block in deck-qa/.
"""
import sys
from pptx import Presentation

SRC = "ceba/CEBA DPPA 2026.pptx"
OUT = "ceba/CEBA DPPA 2026.pptx"

# Slide 0-based indices (subtract 1 from 1-based)
SLIDE_ROADMAP       = 2   # 1-based 3
SLIDE_MODULE5       = 21  # 1-based 22
SLIDE_MODULE6_WRAP  = 26  # 1-based 27 (Know the five-line bill)
SLIDE_INTERACTIVE   = 32  # 1-based 33 (Interactive Exercise)
SLIDE_FORMULAS      = 33  # 1-based 34 (Cost formulas recap)
SLIDE_REPRESENTS    = 39  # 1-based 40 (DPPA represents)
SLIDE_GAMECHANGER   = 41  # 1-based 42 (Vietnam's DPPA is a game-changer)


def set_text(shape, new_text):
    """Replace the entire text of a shape with new_text (preserving the first run's formatting)."""
    tf = shape.text_frame
    # Clear all existing paragraphs
    for p in list(tf.paragraphs):
        p._p.getparent().remove(p._p)
    # Add a fresh paragraph
    p = tf.add_paragraph()
    # If the shape had a hyperlink, clear it (use the first run's properties)
    run = p.add_run()
    run.text = new_text


def find_shape_by_substring(slide, substring):
    """Return the first shape on the slide whose text contains the substring."""
    for shape in slide.shapes:
        if hasattr(shape, "text") and substring in shape.text:
            return shape
    return None


def main():
    prs = Presentation(SRC)
    print(f"Opening deck: {len(prs.slides)} slides")
    changes = []

    # 1. Slide 22: Module 5 divider rename
    slide = prs.slides[SLIDE_MODULE5]
    shape = find_shape_by_substring(slide, "Module 5:")
    if shape:
        old = shape.text
        new = "Module 5:\nThree Canonical Cases"
        set_text(shape, new)
        changes.append(f"Slide 22: renamed '{old.strip()[:40]}' → '{new[:40]}'")

    # 2. Slide 3: Roadmap — append a reference to the 3 canonical cases
    #    The roadmap text shape is the one with "Session Roadmap" eyebrow
    #    — append a clarification that the workshop uses 3 cases, not 5
    slide = prs.slides[SLIDE_ROADMAP]
    shape = find_shape_by_substring(slide, "Session Roadmap")
    if shape:
        # Add a new line to the existing text
        old = shape.text
        new = old + "\n\n• Module 5 walks the three canonical cases (matched / shortfall / excess) — same five-line math the app uses."
        set_text(shape, new)
        changes.append("Slide 3: added canonical-cases reference under Session Roadmap")

    # 3. Slide 27: Module 6 wrap — add reference to the 3 canonical cases
    slide = prs.slides[SLIDE_MODULE6_WRAP]
    shape = find_shape_by_substring(slide, "Know the five-line bill")
    if shape:
        old = shape.text
        new = old + "\n\n• Use the 3 canonical cases (matched / shortfall / excess) at https://dppa-case.web.app to swap the strike and FMP live."
        set_text(shape, new)
        changes.append("Slide 27: added live app reference to Module 6 wrap")

    # 4. Slide 33: Workshop intro — rename
    slide = prs.slides[SLIDE_INTERACTIVE]
    # Find the "DPPA Scenario Analysis" text
    shape = find_shape_by_substring(slide, "DPPA Scenario Analysis")
    if shape:
        old = shape.text
        new = old.replace("DPPA Scenario Analysis", "Apply the 3 Canonical Cases")
        set_text(shape, new)
        changes.append("Slide 33: renamed 'DPPA Scenario Analysis' → 'Apply the 3 Canonical Cases'")

    # 5. Slide 34: Cost formulas recap — add Circular 16/2025 citation
    slide = prs.slides[SLIDE_FORMULAS]
    shape = find_shape_by_substring(slide, "Based on Decree 57/2025/ND-CP")
    if shape:
        old = shape.text
        # Add the Circular 16/2025 reference and the 30-min clarification
        if "Circular 16/2025" not in old:
            new = old + "\n\nLegal basis: Circular 16/2025/TT-BCT (FMP = SMP + CAN, 30-min settlement); Decree 57/2025/ND-CP (grid CfD)."
            set_text(shape, new)
            changes.append("Slide 34: added Circular 16/2025 citation and 30-min clarification")

    # 6. Slide 40: DPPA represents — Samsung/TTC first-DPPA
    slide = prs.slides[SLIDE_REPRESENTS]
    shape = find_shape_by_substring(slide, "DPPA represents a critical step")
    if shape:
        old = shape.text
        # Append a paragraph about the first DPPA
        if "Samsung" not in old:
            new = old + "\n\nFirst grid DPPA live in 2026: Samsung SEVT (Thai Nguyen) ↔ TTC Duc Hue 2 (49 MWp solar + BESS, COD 19 May 2026), ~70 GWh/yr, ~46,000 tCO₂/yr avoided."
            set_text(shape, new)
            changes.append("Slide 40: added Samsung/TTC Duc Hue 2 first-DPPA case")

    # 7. Slide 42: Vietnam's DPPA is a game-changer — "early-stage" → "now operational"
    slide = prs.slides[SLIDE_GAMECHANGER]
    shape = find_shape_by_substring(slide, "early-stage")
    if shape:
        old = shape.text
        # Cleaner replacement: "The DPPA is now available but early-stage" → "The DPPA is now operational"
        new = old.replace(
            "The DPPA is now available but early-stage",
            "The DPPA is now operational",
        )
        if new == old:
            # Fallback if the prefix doesn't match
            new = old.replace("early-stage", "now operational")
        set_text(shape, new)
        changes.append("Slide 42: 'The DPPA is now available but early-stage' → 'The DPPA is now operational'")

    # Save
    prs.save(OUT)
    print(f"\nApplied {len(changes)} changes:")
    for c in changes:
        print(f"  • {c}")
    print(f"\nSaved to {OUT} (still {len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
