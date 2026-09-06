# LIVE: run by CI's deck-parity job (.github/workflows/ci.yml). Regenerate: PYTHONPATH= py verify_deck_numbers.py
"""PHASE-04 (October readiness hardening plan; extended PHASE-04 of
plans/2026-08-22-delivery-stall-recovery-plan.md, 2026-08-23): reconciles every
grouped VND-millions figure shown on a slide BODY *and* in its speaker notes
against the canonical assets/teaching/spine-s1.json and assets/teaching/gate-sweep.json
exports, so a future scenario or engine change cannot silently leave a stale
number on a slide or in the words a presenter reads aloud. Mirrors ASM-001 in
the Oct 2026 teaching-revamp plan: no slide figure is hand-typed without
reconciling to the settlement engine's exports.

Notes were excluded in the original version on the theory that they
"intentionally carry exact answer-key numbers" -- true, but that made notes
the least-verified numbers in the whole pipeline, and they are exactly the
numbers a presenter says out loud to a room that may contain lenders.

--lang {en,vi,zh} selects the expected thousands-grouping character for the
target deck's typography (en/zh: ",", vi: "."; see
app/src/modules/formatters.js's LOCALE_BY_LANG, which this mirrors) and
--deck PATH points at a specific built .pptx, so this can run against a
translated build once one exists. Comparison is done on the normalized
(grouping-stripped) digit string, so "11,020" and "11.020" are recognized as
the same figure regardless of which language produced them.

Run: PYTHONPATH= py verify_deck_numbers.py
     PYTHONPATH= py verify_deck_numbers.py --lang vi --deck "ceba/DPPA Presentation Oct 2026 To Teach vi.pptx"
Exit 0 + "PARITY PASS" on success; exit 1 + offending (slide, origin, token)
triples on mismatch.
"""
import argparse
import json
import os
import re
import sys

from pptx import Presentation

DEFAULT_DECK = os.path.join("ceba", "DPPA Presentation Oct 2026 To Teach.pptx")
ASSETS = os.path.join("assets", "teaching")

# Matches app/src/modules/formatters.js's LOCALE_BY_LANG: en/zh group thousands
# with ",", vi groups with ".".
GROUPING_SEPARATOR_BY_LANG = {"en": ",", "vi": ".", "zh": ","}

# Pedagogical constants shown on-slide that are not raw spine/sweep export
# fields. Each entry must carry a comment explaining exactly where it comes
# from, so this set cannot silently grow into an escape hatch for typos.
# Stored normalized (no grouping separator) since comparison is now locale-
# agnostic.
EXTRA_ALLOWED = {
    # M2 Sankey slide body: "fees {systemService + diffClearing}" — the sum of
    # two spine line items (1,800 + 817), not a field of its own.
    "2617",
}


def number_pattern_for_lang(lang):
    """Return the compiled regex matching a grouped number in `lang`'s typography."""
    sep = re.escape(GROUPING_SEPARATOR_BY_LANG.get(lang, ","))
    return re.compile(rf"\d{{1,3}}(?:{sep}\d{{3}})+")


def normalize_token(token, lang="en"):
    """Strip `lang`'s grouping separator, returning the bare digit string."""
    sep = GROUPING_SEPARATOR_BY_LANG.get(lang, ",")
    return token.replace(sep, "")


def collect_spine_numbers(spine):
    numbers = set()

    # vndMillionsRounded fields are the deck's VND-figure vocabulary; kWh-suffixed
    # fields are its volume vocabulary (e.g. the 5,000,000 kWh monthly load
    # spoken in speaker notes) -- both are genuine settlement-engine outputs,
    # just in different units. Extending to "kwh" here (2026-08-23) is what
    # surfaced slide 16's notes correctly reconciling to /factory/monthlyLoadKwh
    # once notes were added to the scan.
    def walk(node):
        if isinstance(node, dict):
            for key, value in node.items():
                key_lower = key.lower()
                if key_lower.endswith(("vndmillionsrounded", "kwh")) and isinstance(value, (int, float)):
                    numbers.add(str(round(value)))
                else:
                    walk(value)
        elif isinstance(node, list):
            for item in node:
                walk(item)

    walk(spine)
    return numbers


def collect_sweep_numbers(sweep):
    numbers = set()
    for strike in sweep.get("strikes", []):
        numbers.add(str(round(strike)))
    numbers.add(str(sweep.get("passCount", 0)))
    # PHASE-06 (2026-08-23): the per-gate pass counts now cited in the M5
    # notes (see build_oct_teaching_deck.py) alongside the combined count.
    numbers.add(str(sweep.get("buyerPassCount", 0)))
    numbers.add(str(sweep.get("lenderPassCount", 0)))
    numbers.add(str(sweep.get("investorPassCount", 0)))
    numbers.add(str(len(sweep.get("cells", []))))
    return numbers


def allowed_numbers(spine, sweep, extra):
    return collect_spine_numbers(spine) | collect_sweep_numbers(sweep) | extra


def extract_slide_numbers(pptx_path, lang="en"):
    """Return (slide_index, origin, token) triples for every grouped number found
    in each slide's body shapes and, when present, its speaker notes. origin is
    "body" or "notes"."""
    pattern = number_pattern_for_lang(lang)
    prs = Presentation(pptx_path)
    tokens = []
    for slide_index, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for match in pattern.finditer(shape.text_frame.text):
                tokens.append((slide_index, "body", match.group()))
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text
            for match in pattern.finditer(notes_text):
                tokens.append((slide_index, "notes", match.group()))
    return tokens


def main(argv=None):
    if hasattr(sys.stdout, "reconfigure"):
        sys.stdout.reconfigure(encoding="utf-8", errors="replace")
    parser = argparse.ArgumentParser(description=__doc__, formatter_class=argparse.RawDescriptionHelpFormatter)
    parser.add_argument("--deck", default=DEFAULT_DECK, help=f"Path to the built .pptx (default: {DEFAULT_DECK})")
    parser.add_argument(
        "--lang",
        default="en",
        choices=sorted(GROUPING_SEPARATOR_BY_LANG),
        help="Expected thousands-grouping typography of the deck (default: en)",
    )
    args = parser.parse_args(argv)

    if not os.path.exists(args.deck):
        print(f"VERIFY-DECK-NUMBERS: deck not found at {args.deck}", file=sys.stderr)
        return 1

    with open(os.path.join(ASSETS, "spine-s1.json"), encoding="utf-8") as f:
        spine = json.load(f)
    with open(os.path.join(ASSETS, "gate-sweep.json"), encoding="utf-8") as f:
        sweep = json.load(f)

    allowed = allowed_numbers(spine, sweep, EXTRA_ALLOWED)
    tokens = extract_slide_numbers(args.deck, args.lang)

    violations = [
        (slide_index, origin, token)
        for slide_index, origin, token in tokens
        if normalize_token(token, args.lang) not in allowed
    ]

    for slide_index, origin, token in tokens:
        verdict = "OK" if normalize_token(token, args.lang) in allowed else "MISMATCH"
        print(f"slide {slide_index} [{origin}]: {token} -> {verdict}")

    if violations:
        print("\nPARITY FAIL — unreconciled figures found:")
        for slide_index, origin, token in violations:
            print(f"  slide {slide_index} [{origin}]: {token!r} not in spine-s1.json, gate-sweep.json, or EXTRA_ALLOWED")
        return 1

    print(f"\nPARITY PASS — {len(tokens)} figures across {args.deck} all reconcile.")
    return 0


if __name__ == "__main__":
    sys.exit(main(sys.argv[1:]))
