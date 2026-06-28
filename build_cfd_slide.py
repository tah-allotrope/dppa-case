"""Render the consolidated CfD chart AND the three per-scenario CfD charts with the
web-app TOU overlay. For each chart we produce:
  - an animated GIF that sweeps hour-by-hour, highlighting each hour's data points
    with a value callout — plays in Google Slides / PowerPoint present mode
  - an MP4 (H.264 / yuv420p / +faststart, 30 fps) converted from the GIF via ffmpeg,
    for Google Slides video embedding
The consolidated chart additionally produces a static high-res PNG and a 16:9 .pptx.

Outputs (assets/):
  consolidated : cfd-consolidated-chart.png/.gif (+ -vi/-zh-cn gif) and .mp4 variants
  scenarios    : cfd-s{1,2,3}-{en,vi,zh-cn}.gif and matching .mp4

Run with the `py` launcher (matplotlib 3.10, Pillow 12, python-pptx 1.0, ffmpeg on PATH)."""
import io, os, shutil, subprocess
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.lines import Line2D
from PIL import Image

# ---- deck palette ----
TEAL="#0097A7"; AMBER="#FFAB40"; GREEN="#2e9e6b"; MAGENTA="#d6379a"; BLUE="#4285F4"
INK="#212121"; GRAY="#595959"
plt.rcParams["font.family"] = ["Arial","DejaVu Sans"]

# Per-language strings. {strike} is the scenario strike (varies S1=1,250 / S2=1,500 /
# S3=1,250 / consolidated=2,000), formatted with a thousands separator.
TEXTS = {
    "en": {
        "font": "Arial", "off": "Off-peak", "standard": "Standard", "peak": "Peak",
        "load": "Factory load", "solar": "Solar generation", "matched": "Matched volume",
        "strike": "Strike (dotted, VND/kWh)",
        "main": "DPPA settlement: the overlap and the price clamp",
        "hour": "{h}:00   load {load} · solar {solar} · matched {matched} kWh",
        "below": "FMP {fmp} < strike {strike}  →  you top up developer",
        "above": "FMP {fmp} > strike {strike}  →  developer pays you",
        "equal": "FMP = strike  →  no CfD",
    },
    "vi": {
        "font": "Arial", "off": "Thấp điểm", "standard": "Bình thường", "peak": "Cao điểm",
        "load": "Phụ tải nhà máy", "solar": "Điện mặt trời", "matched": "Sản lượng khớp",
        "strike": "Giá thực hiện (nét chấm, VND/kWh)",
        "main": "Thanh toán DPPA: sản lượng khớp và giá thực hiện",
        "hour": "{h}:00   phụ tải {load} · mặt trời {solar} · khớp {matched} kWh",
        "below": "FMP {fmp} < giá thực hiện {strike}  →  bạn bù cho nhà phát triển",
        "above": "FMP {fmp} > giá thực hiện {strike}  →  nhà phát triển trả bạn",
        "equal": "FMP = giá thực hiện  →  không phát sinh CfD",
    },
    "zh": {
        "font": "Microsoft YaHei", "off": "谷时段", "standard": "平时段", "peak": "峰时段",
        "load": "工厂负荷", "solar": "光伏发电", "matched": "匹配电量",
        "strike": "执行价（点线，VND/kWh）",
        "main": "DPPA结算：匹配电量与执行价",
        "hour": "{h}:00   负荷 {load} · 光伏 {solar} · 匹配 {matched} kWh",
        "below": "FMP {fmp} < 执行价 {strike}  →  您向开发商补差价",
        "above": "FMP {fmp} > 执行价 {strike}  →  开发商向您支付差价",
        "equal": "FMP = 执行价  →  无CfD差额",
    },
}

hours = list(range(24))

# ---- Vietnam TOU bands, Mon-Sat (Decision 963/QD-BCT, effective 22 Apr 2026) ----
# Sunday has no peak period: 00:00-06:00 off-peak; 06:00-24:00 standard.
OFF=("#47d7ff",0.10,"#1f8fb0"); STD=("#ffd84f",0.13,"#b08900"); PK=("#ff68d8",0.10,"#c0379a")
BANDS=[(0,6,"off",OFF),(6,17.5,"standard",STD),
       (17.5,22.5,"peak",PK),(22.5,24,"standard",STD)]

# ============================================================
#  SCENARIOS — the consolidated illustrative chart plus the three
#  workshop teaching cases. Daily kWh/FMP curves are illustrative shapes
#  consistent with each scenario's story; the strike line + FMP curve carry
#  the price-clamp teaching. Monthly settlement totals live in the lessons
#  (research/2026-06-29_dppa-scenario-numbers-spec.md), not these daily curves.
# ============================================================
def _matched(load, solar):
    return [min(a, b) for a, b in zip(load, solar)]

CONSOLIDATED = {
    "code": None,  # no S-code → preserves the original consolidated title/asset
    "stub": "cfd-consolidated-chart",
    "strike": 2000,
    "load":  [3000,3000,3000,3000,3000,3000,4000,4000,4000,4700,4700,4700,4700,4700,4700,4700,3900,3900,3900,3900,3200,3200,3200,3200],
    "solar": [0,0,0,0,0,0,0,833,1935,2999,3887,4493,4700,4493,3887,2999,1935,833,0,0,0,0,0,0],
    "fmp":   [1190,1173,1156,1173,1224,1326,1428,1496,1564,1649,1700,1768,1836,1887,1955,2006,2074,2176,2312,2414,2210,1836,1564,1360],
    "kwh_ylim": (0, 5300),
    "fmp_ylim": (800, 2600),
    "fmp_ticks": [1000, 1500, 2000, 2500],
    "fmp_ticklabels": ["1.0k", "1.5k", "2.0k", "2.5k"],
}

# S1 Matched — load ≈ solar overlap; FMP entirely below strike 1,250 (factory tops up).
S1 = {
    "code": "S1", "stub": "cfd-s1", "strike": 1250,
    "load":  [3200,3200,3200,3200,3200,3400,3800,4200,4500,4700,4700,4700,4700,4700,4700,4600,4300,4000,3800,3600,3400,3300,3200,3200],
    "solar": [0,0,0,0,0,0,300,1500,2800,3800,4400,4650,4700,4650,4400,3800,2800,1500,300,0,0,0,0,0],
    "fmp":   [1180,1170,1160,1165,1180,1200,1210,1190,1150,1110,1080,1060,1050,1060,1080,1110,1150,1190,1220,1235,1240,1220,1200,1190],
    "kwh_ylim": (0, 5300),
    "fmp_ylim": (950, 1450),
    "fmp_ticks": [1000, 1100, 1200, 1300, 1400],
    "fmp_ticklabels": ["1.0k", "1.1k", "1.2k", "1.3k", "1.4k"],
}

# S2 Shortfall — load clearly above solar; FMP entirely above strike 1,500 (developer pays you).
S2 = {
    "code": "S2", "stub": "cfd-s2", "strike": 1500,
    "load":  [4500,4500,4500,4500,4600,4700,4800,4900,5000,5000,5000,5000,5000,5000,5000,4900,4800,4700,4600,4600,4500,4500,4500,4500],
    "solar": [0,0,0,0,0,0,200,900,1700,2300,2750,2950,3000,2950,2750,2300,1700,900,200,0,0,0,0,0],
    "fmp":   [1560,1545,1535,1545,1575,1620,1660,1690,1700,1690,1660,1640,1630,1640,1670,1710,1740,1770,1760,1730,1690,1650,1610,1580],
    "kwh_ylim": (0, 5300),
    "fmp_ylim": (1300, 1950),
    "fmp_ticks": [1400, 1500, 1600, 1700, 1800],
    "fmp_ticklabels": ["1.4k", "1.5k", "1.6k", "1.7k", "1.8k"],
}

# S3 Excess — solar clearly above load midday (over-generation); sunny-month FMP dips
# well below strike 1,250. The excess (solar above load) settles nothing — spot only.
S3 = {
    "code": "S3", "stub": "cfd-s3", "strike": 1250,
    "load":  [3600,3600,3600,3600,3600,3700,3800,3900,4000,4000,4000,4000,4000,4000,4000,3900,3800,3700,3600,3600,3600,3600,3600,3600],
    "solar": [0,0,0,0,0,0,400,1800,3200,4200,4800,5050,5100,5050,4800,4200,3200,1800,400,0,0,0,0,0],
    "fmp":   [1190,1180,1170,1175,1190,1200,1180,1130,1070,1010,970,940,930,940,970,1010,1070,1130,1180,1210,1220,1210,1200,1190],
    "kwh_ylim": (0, 5300),
    "fmp_ylim": (850, 1450),
    "fmp_ticks": [900, 1050, 1200, 1350],
    "fmp_ticklabels": ["0.9k", "1.05k", "1.2k", "1.35k"],
}

for _s in (CONSOLIDATED, S1, S2, S3):
    _s["matched"] = _matched(_s["load"], _s["solar"])


def draw_base(axL, scen, lang="en"):
    t = TEXTS[lang]
    load, solar, matched, fmp = scen["load"], scen["solar"], scen["matched"], scen["fmp"]
    kwh_top = scen["kwh_ylim"][1]
    axR = axL.twinx()
    # TOU overlay (behind everything)
    for s,e,name,(fill,a,lc) in BANDS:
        axL.axvspan(s-0.5, e-0.5, color=fill, alpha=a, zorder=0, lw=0)
        axL.text((s+e)/2-0.5, kwh_top-220, t[name], ha="center", va="top", fontsize=9.5,
                 color=lc, fontweight="bold", zorder=1, fontfamily=t["font"])
    # left axis (kWh/h)
    axL.fill_between(hours, load,  color=TEAL,  alpha=0.10, zorder=1)
    axL.plot(hours, load,  color=TEAL,  lw=3.0, label=t["load"], zorder=2)
    axL.plot(hours, solar, color=AMBER, lw=3.0, label=t["solar"], zorder=2)
    axL.fill_between(hours, matched, color=GREEN, alpha=0.30, zorder=1)
    axL.plot(hours, matched, color=GREEN, lw=2.0, label=t["matched"], zorder=2)
    axL.set_ylim(*scen["kwh_ylim"]); axL.set_xlim(-0.5, 23.5)
    axL.set_ylabel("kWh / h", color=GRAY, fontsize=13)
    axL.set_xticks(range(0,24,2)); axL.set_xticklabels([f"{h:02d}" for h in range(0,24,2)], color=GRAY, fontsize=11)
    axL.tick_params(axis="y", colors=GRAY, labelsize=11)
    axL.grid(True, axis="y", color="#000000", alpha=0.05, zorder=0)
    for sp in axL.spines.values(): sp.set_color("#E0E0E0")
    # right axis (VND/kWh)
    axR.plot(hours, fmp, color=MAGENTA, lw=3.0, ls=(0,(7,4)), label="FMP (VND/kWh)", zorder=3)
    axR.axhline(scen["strike"], color=BLUE, lw=2.4, ls=(0,(1.2,3.0)),
                dash_capstyle="round", label=t["strike"], zorder=3)
    axR.set_ylim(*scen["fmp_ylim"])
    axR.set_ylabel("VND / kWh", color=MAGENTA, fontsize=13)
    axR.tick_params(axis="y", colors=MAGENTA, labelsize=11)
    axR.set_yticks(scen["fmp_ticks"]); axR.set_yticklabels(scen["fmp_ticklabels"])
    for sp in axR.spines.values(): sp.set_color("#E0E0E0")
    return axR

def legend_below(axL, axR, lang="en"):
    h1,l1 = axL.get_legend_handles_labels(); h2,l2 = axR.get_legend_handles_labels()
    # A long explicit dotted sample remains legible after GIF quantization.
    h2[-1] = Line2D([0],[0], color=BLUE, lw=2.4, linestyle=(0,(1.2,3.0)),
                    dash_capstyle="round")
    axL.legend(h1+h2, l1+l2, loc="upper center", bbox_to_anchor=(0.5,-0.12),
               ncol=5, frameon=False, fontsize=10.5, labelcolor=INK,
               handlelength=3.2, handletextpad=0.7, prop={"family": TEXTS[lang]["font"], "size": 10.5})

SUB=("Load & solar overlap = matched (settles on CfD)  ·  TOU: Mon-Sat, Decision 963 (22 Apr 2026)  ·  "
     "FMP below strike → you top up; above → developer pays you")

def fmt(n): return f"{n:,}"

def convert_to_mp4(gif_path):
    """GIF → MP4 (H.264 / yuv420p / +faststart, 30 fps) for Google Slides."""
    if not shutil.which("ffmpeg"):
        print("  ffmpeg not on PATH — skipping MP4 for", gif_path); return
    mp4 = os.path.splitext(gif_path)[0] + ".mp4"
    cmd = ["ffmpeg","-y","-i",gif_path,"-movflags","+faststart","-pix_fmt","yuv420p",
           "-vf","fps=30,scale=trunc(iw/2)*2:trunc(ih/2)*2","-c:v","libx264",mp4]
    subprocess.run(cmd, check=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
    print("  MP4:", mp4, os.path.getsize(mp4), "bytes")

# ---------- animated GIF (hour-by-hour highlight) ----------
def generate_gif(scen, lang, filename, make_mp4=True):
    t = TEXTS[lang]
    load, solar, matched, fmp, strike = scen["load"], scen["solar"], scen["matched"], scen["fmp"], scen["strike"]
    suptitle = (f"{scen['code']} · {t['main']}" if scen.get("code") else t["main"])
    frames=[]
    for h in range(24):
        fig, axL = plt.subplots(figsize=(11.0,6.2), dpi=100)
        fig.patch.set_facecolor("white"); axL.set_facecolor("white")
        axR = draw_base(axL, scen, lang); legend_below(axL, axR, lang)
        # vertical guide + highlighted markers
        axL.axvline(h, color="#212121", alpha=0.30, lw=1.2, zorder=4)
        for val,c in ((load[h],TEAL),(solar[h],AMBER),(matched[h],GREEN)):
            axL.scatter([h],[val], s=120, color=c, edgecolor="white", linewidth=1.6, zorder=6)
        axR.scatter([h],[fmp[h]], s=120, color=MAGENTA, edgecolor="white", linewidth=1.6, zorder=6)
        if fmp[h] < strike: dirn=t["below"].format(fmp=fmt(fmp[h]), strike=fmt(strike))
        elif fmp[h] > strike: dirn=t["above"].format(fmp=fmt(fmp[h]), strike=fmt(strike))
        else: dirn=t["equal"]
        title=t["hour"].format(h=f"{h:02d}", load=fmt(load[h]), solar=fmt(solar[h]), matched=fmt(matched[h])) + "\n" + dirn
        fig.suptitle(suptitle, x=0.5, y=0.985, fontsize=15, fontweight="bold",
                     color="#00727e", fontfamily=t["font"])
        axL.set_title(title, fontsize=11, color=INK, pad=8, fontfamily=t["font"])
        fig.subplots_adjust(left=0.07, right=0.93, top=0.84, bottom=0.18)
        buf=io.BytesIO(); fig.savefig(buf, format="png", facecolor="white"); plt.close(fig)
        buf.seek(0); frames.append(Image.open(buf).convert("RGB"))
    pframes=[f.quantize(colors=128, method=Image.MEDIANCUT) for f in frames]
    gif=os.path.join("assets", filename)
    pframes[0].save(gif, save_all=True, append_images=pframes[1:], duration=480,
                    loop=0, optimize=True, disposal=2)
    print("GIF:", gif, os.path.getsize(gif), "bytes,", len(pframes), "frames")
    if make_mp4: convert_to_mp4(gif)

# ---------- consolidated static PNG ----------
fig, axL = plt.subplots(figsize=(12.8,6.8), dpi=200)
fig.patch.set_facecolor("white"); axL.set_facecolor("white")
axR = draw_base(axL, CONSOLIDATED); legend_below(axL, axR)
fig.suptitle("DPPA settlement: the overlap and the price clamp", x=0.5, y=0.98,
             fontsize=18, fontweight="bold", color="#00727e")
axL.set_title(SUB, fontsize=10.5, color=GRAY, pad=8)
fig.subplots_adjust(left=0.07, right=0.93, top=0.88, bottom=0.17)
png = os.path.join("assets","cfd-consolidated-chart.png")
fig.savefig(png, dpi=200, facecolor="white"); plt.close(fig)
print("PNG:", png, os.path.getsize(png), "bytes")

# ---------- consolidated animated GIFs + MP4s (en/vi/zh) ----------
# The en consolidated GIF keeps its historical name (no -en); its MP4 is created
# explicitly below as -en.mp4, so skip the auto-MP4 to avoid a redundant file.
generate_gif(CONSOLIDATED, "en", "cfd-consolidated-chart.gif", make_mp4=False)
generate_gif(CONSOLIDATED, "vi", "cfd-consolidated-chart-vi.gif")
generate_gif(CONSOLIDATED, "zh", "cfd-consolidated-chart-zh-cn.gif")
# the en consolidated MP4 keeps its historical name (no -en on the gif)
if os.path.exists(os.path.join("assets","cfd-consolidated-chart.gif")):
    src = os.path.join("assets","cfd-consolidated-chart.gif")
    dst = os.path.join("assets","cfd-consolidated-chart-en.mp4")
    if shutil.which("ffmpeg"):
        subprocess.run(["ffmpeg","-y","-i",src,"-movflags","+faststart","-pix_fmt","yuv420p",
                        "-vf","fps=30,scale=trunc(iw/2)*2:trunc(ih/2)*2","-c:v","libx264",dst],
                       check=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
        print("  MP4:", dst, os.path.getsize(dst), "bytes")

# ---------- per-scenario animated GIFs + MP4s (en/vi/zh) ----------
LANGS = [("en","en"), ("vi","vi"), ("zh","zh-cn")]
for scen in (S1, S2, S3):
    for lang_key, lang_suffix in LANGS:
        generate_gif(scen, lang_key, f"{scen['stub']}-{lang_suffix}.gif")

# ---------- 16:9 slide .pptx (static consolidated PNG) ----------
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
prs=Presentation(); prs.slide_width=Inches(10); prs.slide_height=Inches(5.625)
slide=prs.slides.add_slide(prs.slide_layouts[6])
def tb(l,t,w,h,text,size,color,bold=True):
    b=slide.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h)); f=b.text_frame; f.word_wrap=True
    p=f.paragraphs[0]; r=p.add_run(); r.text=text
    r.font.size=Pt(size); r.font.bold=bold; r.font.name="Arial"; r.font.color.rgb=RGBColor.from_string(color)
tb(0.45,0.22,9.1,0.5,"DPPA settlement: the overlap and the price clamp",22,"00727E")
tb(0.45,0.70,9.1,0.4,"New TOU bands (Decision 963) + load/solar/matched + FMP vs strike",12,"595959",bold=False)
iw,ih=Image.open(png).size; disp_w=9.0; disp_h=disp_w*ih/iw; top=1.18
if top+disp_h>5.45: disp_h=5.45-top; disp_w=disp_h*iw/ih
slide.shapes.add_picture(png, Inches((10-disp_w)/2), Inches(top), width=Inches(disp_w))
tb(0.45,5.34,9.1,0.3,"Illustrative balanced day · strike 2,000 · FMP ~1,700 avg (NSMO/ERAV not public) · mirrors dppa-case.web.app",9,"78909C",bold=False)
out=os.path.join("ceba","DPPA CfD Consolidated Chart.pptx"); prs.save(out)
print("PPTX:", out, os.path.getsize(out), "bytes")
