"""PHASE-03 (Modules 1-6 teaching revamp): build the October 2026 teaching deck
from the 44-slide master `ceba/CEBA DPPA 2026.pptx`, following the content spec
in TEACH_SLIDES below (TASK-03-01). Visual-first, <=30 words/content slide,
plain-language-first (Decree-57 symbols deferred to the M6 decoder), one
hidden fallback slide after every divider, full speaker notes.

Run: PYTHONPATH= py build_oct_teaching_deck.py --lang en
Outputs: ceba/DPPA Presentation Oct 2026 To Teach {lang}.pptx (en has no suffix)
"""
import argparse, copy, json, os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

MASTER = os.path.join("ceba", "CEBA DPPA 2026.pptx")
ASSETS = os.path.join("assets", "teaching")
TEAL = "00727E"
INK = "212121"
GRAY = "595959"

with open(os.path.join(ASSETS, "spine-s1.json"), encoding="utf-8") as f:
    SPINE = json.load(f)

L = SPINE["bill"]["lines"]
CEVN = SPINE["bill"]["cEvn"]["vndMillionsRounded"]
CKH = SPINE["bill"]["cKh"]["vndMillionsRounded"]
BAU = SPINE["comparison"]["bauMonthlyVndMillionsRounded"]
FACTORY = SPINE["factory"]["name"]

TEXT = {
    "en": {
        "cold_open_title": f"{FACTORY}: one month, two bills",
        "cold_open_body": f"Today: {BAU:,} tr VND. With a DPPA: {CKH:,} tr VND. Where did the gap come from? You will compute it.",
        "divider": ["1. The Baseline", "2. The Bill", "3. The Lock", "4. Three Doors", "5. The Case", "6. Decide"],
        "checkpoint": [
            "If load rises in the evening peak, does your bill rise faster or slower than at noon?",
            "Which of the five lines disappears when consumption exactly equals matched volume?",
            "If the market price jumps above the strike mid-afternoon, who pays whom?",
            "Which of the three doors is hardest to pass when the strike is set too high?",
            "In the 56-scenario sweep, how many pass all three gates at once?",
            "Name one lever that would have flipped this month's CfD sign.",
        ],
        "m1_title": "What you pay EVN today",
        "m1_body": f"{FACTORY} — hour by hour, before any DPPA. This is the bill every DPPA offer is judged against.",
        "m2a_title": "How much energy actually settles?",
        "m2a_body": "Generation leaks a little, then passes a load gate and a contract gate. Only what survives both gates settles.",
        "m2b_title": "The five-line bill, one arrow at a time",
        "m2b_body": f"Market energy {L['marketEnergy']['vndMillionsRounded']:,} + fees {L['systemService']['vndMillionsRounded']+L['diffClearing']['vndMillionsRounded']:,} + CfD {L['cfd']['vndMillionsRounded']:,} = {CKH:,} tr VND.",
        "m3_title": "A lock, not a discount",
        "m3_body": "The strike is a seesaw. Below it, you top up the developer. Above it, the developer pays you.",
        "m3_app": "App moment: drag the market-price slider through the strike and watch the CfD line flip sign.",
        "m4_title": "Three doors the deal must pass",
        "m4_body": "Buyer: cost below doing nothing. Lender: covered every year. Investor: return earned. All three, or no deal.",
        "m5_title": "56 scenarios, one empty window",
        "m5_body": "Sweep strike x volume across both case studies: zero of 56 combinations clear all three doors at once.",
        "m5_exercise": "Your turn: compute this factory's five-line bill from the worksheet, then verify it in the app.",
        "m6_decoder_title": "Decoder: your words, the decree's symbols",
        "m6_decoder_rows": [
            ("Matched / settled volume", "Q_khc"),
            ("Total consumption", "Q_KH"),
            ("Market price", "FMP"),
            ("Locked price", "Strike / P_c"),
            ("Total bill", "C_KH = C_EVN + C_CfD"),
        ],
        "m6_levers_title": "Five levers you can still negotiate",
        "m6_levers": SPINE["levers"],
        "close_title": "You can now compute every line of this bill",
        "close_body": "Bring the five levers to your next term sheet. Scan for the app, lessons, and the 90-min scenario workshop.",
        "appendix_titles": ["Scenario 1 — Matched", "Scenario 2 — Shortfall", "Scenario 3 — Excess"],
        "appendix_takeaways": [
            "Consumption equals matched volume. No residual retail line.",
            "Consumption exceeds contract. The residual buys back at retail — always the expensive line.",
            "Generation exceeds consumption. The excess earns spot revenue only — no CfD, no free upside.",
        ],
    },
}
# vi / zh-cn text layers are added after EN content freeze (PHASE-06); the
# script already supports --lang for when those layers are written.
TEXT["vi"] = TEXT["en"]
TEXT["zh"] = TEXT["en"]


def word_count(*parts):
    return sum(len(p.split()) for p in parts if p)


def add_textbox(slide, left, top, width, height, text, size, color=INK, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = "Arial"
    run.font.color.rgb = RGBColor.from_string(color)
    return box


def add_picture_fit(slide, path, left, top, max_width, max_height):
    from PIL import Image
    iw, ih = Image.open(path).size
    ratio = min(max_width / iw, max_height / ih)
    w, h = iw * ratio, ih * ratio
    x = left + (max_width - w) / 2
    y = top + (max_height - h) / 2
    slide.shapes.add_picture(path, Inches(x), Inches(y), width=Inches(w), height=Inches(h))


def set_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def hide_slide(slide):
    slide._element.set("show", "0")


def clear_slides(prs):
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    for sld in slides:
        rId = sld.get(qn("r:id"))
        prs.part.drop_rel(rId)
        xml_slides.remove(sld)


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[9])  # BLANK


def cold_open(prs, t):
    s = blank_slide(prs)
    add_textbox(s, 0.5, 0.25, 9.0, 0.6, t["cold_open_title"], 24, TEAL, bold=True)
    add_picture_fit(s, os.path.join(ASSETS, "cold-open-bill-pair-en.png"), 0.5, 0.9, 9.0, 3.6)
    add_textbox(s, 0.5, 4.6, 9.0, 0.4, t["cold_open_body"], 12, GRAY)
    set_notes(s, "COLD OPEN (2 min). Show the two bars. Do not explain the gap yet — that is the session's promise. "
                 "Hook line: 'In 60 minutes you will compute every line of that difference yourself.' "
                 "Breadcrumb strip introduced next as the roadmap.")
    assert word_count(t["cold_open_title"], t["cold_open_body"]) <= 30
    return s


def divider(prs, module_num, t):
    s = blank_slide(prs)
    add_textbox(s, 0.5, 1.8, 9.0, 0.8, t["divider"][module_num - 1], 30, TEAL, bold=True, align=PP_ALIGN.CENTER)
    breadcrumb = os.path.join(ASSETS, f"breadcrumb-strip-m{module_num}-en.png")
    add_picture_fit(s, breadcrumb, 0.3, 3.0, 9.4, 1.3)
    set_notes(s, f"CHECKPOINT (30 sec, show of hands): {t['checkpoint'][module_num - 1]} "
                 f"Do not move on until at least half the room answers.")
    return s


def fallback_slide(prs, module_num, gif_path=None, note=""):
    from PIL import Image
    s = blank_slide(prs)
    add_textbox(s, 0.5, 0.3, 9.0, 0.5, f"[Fallback — Module {module_num} app demo]", 16, GRAY, bold=True)
    mp4_path = os.path.join(ASSETS, "fallback", f"teach-m{module_num}.mp4")
    poster_path = os.path.join(ASSETS, "fallback", f"teach-m{module_num}-poster.png")
    if os.path.exists(mp4_path):
        iw, ih = Image.open(poster_path).size if os.path.exists(poster_path) else (1280, 720)
        left, top, max_width, max_height = 0.5, 0.9, 9.0, 3.9
        ratio = min(max_width / iw, max_height / ih)
        w_in, h_in = iw * ratio, ih * ratio
        x_in = left + (max_width - w_in) / 2
        y_in = top + (max_height - h_in) / 2
        s.shapes.add_movie(
            mp4_path, Inches(x_in), Inches(y_in), Inches(w_in), Inches(h_in),
            poster_frame_image=poster_path if os.path.exists(poster_path) else None,
            mime_type="video/mp4",
        )
    elif gif_path and os.path.exists(gif_path):
        add_picture_fit(s, gif_path, 0.5, 0.9, 9.0, 3.9)
    else:
        add_textbox(s, 0.5, 1.5, 9.0, 1.0,
                    "Recorded demo not yet captured. Run: npm run record:demos (app/) — "
                    "this regenerates assets/teaching/fallback/teach-m*.mp4 from the live ?teach=1 steps.",
                    12, GRAY)
    if note:
        add_textbox(s, 0.5, 4.9, 9.0, 0.3, note, 9, GRAY)
    hide_slide(s)
    set_notes(s, f"HIDDEN FALLBACK SLIDE. Unhide and play only if the live app fails for the module {module_num} demo. "
                  "Also runnable locally via `npm run preview` in app/ if venue wifi is down.")
    return s


def content_slide(prs, title, body, image_path, notes, checklist=None):
    s = blank_slide(prs)
    add_textbox(s, 0.5, 0.25, 9.0, 0.5, title, 22, TEAL, bold=True)
    add_picture_fit(s, image_path, 0.5, 0.85, 9.0, 3.7)
    add_textbox(s, 0.5, 4.6, 9.0, 0.4, body, 11.5, GRAY)
    set_notes(s, notes)
    assert word_count(title, body) <= 30, f"Slide '{title}' exceeds 30-word budget: {word_count(title, body)}"
    return s


def m6_decoder_slide(prs, t):
    s = blank_slide(prs)
    add_textbox(s, 0.5, 0.25, 9.0, 0.5, t["m6_decoder_title"], 22, TEAL, bold=True)
    y = 1.0
    for plain, symbol in t["m6_decoder_rows"]:
        add_textbox(s, 0.6, y, 4.6, 0.4, plain, 13, INK)
        add_textbox(s, 5.4, y, 3.8, 0.4, symbol, 13, TEAL, bold=True)
        y += 0.55
    set_notes(s, "This is the ONLY slide with Decree-57 notation. Frame it as translation, not new content: "
                 "'you already understand every row on the left — the right column is what the legal text calls it.' "
                 "Full formula appendix (old slide 27) is in the printed A4 card, not here.")
    return s


def m6_levers_slide(prs, t):
    s = blank_slide(prs)
    add_textbox(s, 0.5, 0.25, 9.0, 0.5, t["m6_levers_title"], 22, TEAL, bold=True)
    y = 1.0
    for lever in t["m6_levers"]:
        add_textbox(s, 0.6, y, 8.8, 0.6, f"{lever['name']}: {lever['movesWhen']}", 12, INK)
        y += 0.75
    set_notes(s, "Walk each lever in ~15 seconds: name it, say which way it moves the bill. "
                 "This is the checklist they take into their next term-sheet negotiation.")
    return s


def close_slide(prs, t):
    s = blank_slide(prs)
    add_textbox(s, 0.5, 0.4, 9.0, 0.6, t["close_title"], 22, TEAL, bold=True, align=PP_ALIGN.CENTER)
    add_picture_fit(s, os.path.join(ASSETS, "cold-open-bill-pair-en.png"), 1.0, 1.1, 8.0, 3.0)
    add_textbox(s, 0.5, 4.3, 9.0, 0.6, t["close_body"], 12, GRAY, align=PP_ALIGN.CENTER)
    set_notes(s, "CLOSE (2 min). Point back at the cold-open bars. "
                 "Invite to the 90-min scenario workshop for negotiation practice. QR: dppa-case.web.app and lessons index "
                 "(QR image pending python 'qrcode' package — currently a text URL; see PHASE-03 note).")
    return s


def appendix_slide(prs, idx, title, takeaway, gif_stub):
    s = blank_slide(prs)
    add_textbox(s, 0.5, 0.25, 9.0, 0.5, title, 20, TEAL, bold=True)
    gif_png_fallback = os.path.join("assets", f"{gif_stub}.gif")
    if os.path.exists(gif_png_fallback):
        add_picture_fit(s, gif_png_fallback, 0.5, 0.85, 9.0, 3.7)
    add_textbox(s, 0.5, 4.6, 9.0, 0.4, takeaway, 11.5, GRAY)
    set_notes(s, f"Appendix slide, slimmed per DEC-021. Play the animated chart; state the one-line takeaway. "
                  f"Detailed numbers live in the scenario lessons/worksheets, not here.")
    return s


def build(lang):
    t = TEXT[lang]
    prs = Presentation(MASTER)
    clear_slides(prs)

    cold_open(prs, t)

    content_slide(prs, t["m1_title"], t["m1_body"], os.path.join(ASSETS, "m1-tou-strip-en.png"),
                  "M1 (4 min). Point at the TOU bands, then the load line. State: 'This is BAU — every DPPA offer "
                  "is judged against this.' App moment: none for M1 (visual only).")

    divider(prs, 1, t)
    fallback_slide(prs, 1)

    divider(prs, 2, t)
    content_slide(prs, t["m2a_title"], t["m2a_body"], os.path.join(ASSETS, "m2-funnel-en.png"),
                  "M2a (2 min). Walk the funnel top to bottom: generation, loss leak, load gate, contract gate. "
                  "No formulas — 'you only settle what survives both gates.'")
    content_slide(prs, t["m2b_title"], t["m2b_body"], os.path.join(ASSETS, "m2-sankey-en-5.png"),
                  "M2b (3 min). Play m2-sankey-build-en.gif in slideshow mode (autoplay). Read each arrow aloud as it "
                  "appears. APP MOMENT: switch to the live app, scroll to the five-line-bill panel, confirm the same "
                  f"numbers (C_EVN {CEVN:,}, CfD {L['cfd']['vndMillionsRounded']:,}, C_KH {CKH:,} tr VND).")
    fallback_slide(prs, 2, os.path.join(ASSETS, "m2-sankey-build-en.gif"))

    divider(prs, 3, t)
    content_slide(prs, t["m3_title"], t["m3_body"], os.path.join(ASSETS, "m3-seesaw-en.png"),
                  "M3 (4 min). Show the seesaw, then play assets/cfd-s1-en.gif for the 24h sign flip. " + t["m3_app"])
    fallback_slide(prs, 3, os.path.join("assets", "cfd-s1-en.gif"))

    divider(prs, 4, t)
    content_slide(prs, t["m4_title"], t["m4_body"], os.path.join(ASSETS, "m4-three-doors-en.png"),
                  "M4 (4 min). Name each door once, no ratio math on-slide. APP MOMENT: multi-year panel, point at "
                  "the crossover year as the buyer-door check.")
    fallback_slide(prs, 4)

    divider(prs, 5, t)
    content_slide(prs, t["m5_title"], t["m5_body"], os.path.join(ASSETS, "m5-gate-heatmap-en.png"),
                  "M5 setup (2 min). Do NOT reveal the heatmap takeaway yet — that comes after the exercise.")
    s = blank_slide(prs)
    add_textbox(s, 0.5, 0.3, 9.0, 0.5, "Your turn: compute the bill", 20, TEAL, bold=True)
    add_textbox(s, 0.5, 1.0, 9.0, 3.5,
                f"Worksheet: {FACTORY}, S1. Volumes are pre-filled. Compute all five lines in VND millions, total, "
                "and compare to the BAU bar from the cold open. 10 minutes.", 13, INK)
    hide_none = None
    set_notes(s, "M5 EXERCISE (10 min, the session's success criterion). Hand out the one-page worksheet "
                 "(lessons/0012-reference-card). Volumes pre-filled: contracted=total=5,000,000 kWh. "
                 f"Answer key: marketEnergy {L['marketEnergy']['vndMillionsRounded']:,}, service "
                 f"{L['systemService']['vndMillionsRounded']:,}, clearing {L['diffClearing']['vndMillionsRounded']:,}, "
                 f"residual {L['additionalPurchase']['vndMillionsRounded']:,}, CfD {L['cfd']['vndMillionsRounded']:,}, "
                 f"total {CKH:,} tr VND. APP MOMENT: verify against the five-line-bill panel.")
    fallback_slide(prs, 5)
    content_slide(prs, "The empty window, revealed", "Now scale your month x12 x20 strikes: zero of 56 pass all three doors.",
                  os.path.join(ASSETS, "m5-gate-heatmap-en.png"),
                  "M5 REVEAL (2 min). This is the punchline: the exercise they just did, multiplied across a lifetime "
                  "and a strike sweep, is why the window is empty in these two case studies.")

    divider(prs, 6, t)
    m6_decoder_slide(prs, t)
    m6_levers_slide(prs, t)
    fallback_slide(prs, 6)

    close_slide(prs, t)

    for i, (title, takeaway) in enumerate(zip(t["appendix_titles"], t["appendix_takeaways"])):
        appendix_slide(prs, i + 1, title, takeaway, f"cfd-s{i+1}-en")

    suffix = "" if lang == "en" else f" {lang}"
    out = os.path.join("ceba", f"DPPA Presentation Oct 2026 To Teach{suffix}.pptx")
    prs.save(out)
    print("Saved:", out, "-", len(prs.slides.__iter__.__self__._sldIdLst) if False else len(prs.slides._sldIdLst), "slides")
    return out


if __name__ == "__main__":
    parser = argparse.ArgumentParser()
    parser.add_argument("--lang", default="en", choices=["en", "vi", "zh"])
    args = parser.parse_args()
    build(args.lang)
