"""
Build the 3 canonical case slides for the CEBA DPPA 2026 deck.

Strategy (low-risk):
  1. Open the existing deck
  2. Add 3 new slides at the end using a clean blank layout
  3. Reorder them to position 38 (1-indexed) via XML manipulation of <p:sldIdLst>
     so the new slides become slides 38, 39, 40 and the existing workshop
     scenario slides shift to 41, 42, ...
  4. Delete the old Scenarios 1 (now 41-44) and Scenarios 3 (now 45-48)
  5. Save

The slide content mirrors the app's `walkthroughCaseCard` format so a
facilitator can hand off from deck to live tool without re-doing the math.

Verified 2026 numbers (single source of truth: research/2026-06-22_vietnam-dppa-2026.md):
  - Retail: 2,204.07 VND/kWh
  - Fixed fees: 360 + 163.3 = 523.3 VND/kWh
  - Loss factor: k*K_pp = 1.026 * 1.008 = 1.0342
  - Strike (teaching default): 2,000 VND/kWh
  - FMP: 1,427 VND/kWh (illustrative — NSMO/ERAV not public)

Factory frame: 24h representative day; 4,200 kWh/h load (Case A: balanced);
6,100 load vs 4,200 gen (Case B: shortfall); 2,600 load vs 4,700 gen (Case C: excess).
Matches the app's default-scenarios.js archetypes.
"""
import copy
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

SRC = "ceba/CEBA DPPA 2026.pptx"
OUT = "ceba/CEBA DPPA 2026.pptx"

# Verified 2026 numbers
RETAIL = 2204.07
FEES = 523.3
LOSS = 1.0342  # k * K_pp
STRIKE = 2000
FMP = 1427  # illustrative
KWPP = 1000
HOURS = 24

# Color palette (neon, matches app)
CYAN = RGBColor(0x47, 0xD7, 0xFF)
MINT = RGBColor(0x9A, 0xFF, 0xDE)
AMBER = RGBColor(0xFF, 0xD8, 0x4F)
MAGENTA = RGBColor(0xFF, 0x68, 0xD8)
RED = RGBColor(0xFF, 0x5A, 0x5F)
WHITE = RGBColor(0xF6, 0xFB, 0xFF)
MUTED = RGBColor(0x8B, 0x97, 0xA8)
DARK = RGBColor(0x05, 0x08, 0x16)


def add_text(slide, left, top, width, height, text, *, size=14, bold=False,
             color=WHITE, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """Add a single text box with one line/paragraph."""
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def add_paragraphs(slide, left, top, width, height, lines, *, default_size=12,
                   default_color=WHITE, default_bold=False):
    """Add a textbox with multiple paragraphs.

    Each item in `lines` is either a string (defaults) or a dict:
        {"text": ..., "size": 12, "bold": False, "color": WHITE,
         "align": "left", "bullet": False}
    """
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    for i, item in enumerate(lines):
        if isinstance(item, str):
            item = {"text": item}
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
                       "right": PP_ALIGN.RIGHT}.get(item.get("align", "left"), PP_ALIGN.LEFT)
        run = p.add_run()
        run.text = item["text"]
        run.font.size = Pt(item.get("size", default_size))
        run.font.bold = item.get("bold", default_bold)
        run.font.color.rgb = item.get("color", default_color)
    return tb


def add_filled_rect(slide, left, top, width, height, color):
    """Add a filled rectangle (used for accent strips / backgrounds)."""
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top),
                                    Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_header(slide, eyebrow, title, *, y=0.18):
    """Add the standard eyebrow + title at the top of a content slide."""
    add_text(slide, 0.4, y, 9.2, 0.32, eyebrow, size=11, color=CYAN, bold=True)
    add_text(slide, 0.4, y + 0.30, 9.2, 0.55, title, size=22, color=WHITE, bold=True)


def add_caveat_strip(slide, *, y=5.22):
    """Bottom-of-slide caveat strip (FMP illustrative, etc.)."""
    add_text(slide, 0.4, y, 9.2, 0.22,
             "FMP ~1,427 VND/kWh is illustrative (NSMO/ERAV not publicly published). "
             "Retail 2,204 / fees 523.3 / loss 1.0342 are 2025 verified; "
             "settlement modeled hourly (Circular 16/2025 = 30-min).",
             size=8, color=MUTED)


def add_case_content(slide, *, case_label, header_color, narrative, evn_lines,
                     cfd_lines, net_line, formula_lines, comparison_line, y=0.95):
    """Add the standard canonical-case slide body.

    All y-offsets are relative to the `y` parameter.  The slide is
    5.62" tall (widescreen 16:9 default), so every element must fit
    above the caveat strip at y=5.22.  The comparison punchline must
    land at or before ~5.05" to stay on-slide and leave ~0.17" gap
    before the caveat.
    """
    # Pill: case label
    add_filled_rect(slide, 0.4, y, 1.4, 0.32, header_color)
    add_text(slide, 0.4, y + 0.02, 1.4, 0.28, case_label, size=10, color=DARK,
             bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Narrative
    add_text(slide, 0.4, y + 0.38, 9.2, 0.42, narrative, size=11, color=WHITE)

    # 5-line EVN bill (left column) — compressed 4 lines + total
    add_text(slide, 0.4, y + 0.88, 4.4, 0.24, "EVN bill (5-line, on matched volume)",
             size=9, color=CYAN, bold=True)
    add_paragraphs(slide, 0.4, y + 1.12, 4.4, 1.45, evn_lines,
                   default_size=9, default_color=WHITE)

    # CfD (right column)
    add_text(slide, 5.0, y + 0.88, 4.6, 0.24, "CfD settlement (factory ↔ developer)",
             size=9, color=AMBER, bold=True)
    add_paragraphs(slide, 5.0, y + 1.12, 4.6, 1.45, cfd_lines,
                   default_size=9, default_color=WHITE)

    # Net total (full width)
    add_filled_rect(slide, 0.4, y + 2.68, 9.2, 0.38, DARK)
    add_text(slide, 0.4, y + 2.69, 9.2, 0.36, net_line, size=13, color=MINT,
             bold=True, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)

    # Formula breakdown (2 compact lines)
    add_text(slide, 0.4, y + 3.14, 9.2, 0.24, "Formula breakdown",
             size=9, color=CYAN, bold=True)
    add_paragraphs(slide, 0.4, y + 3.38, 9.2, 0.55, formula_lines,
                   default_size=8, default_color=MUTED)

    # Comparison vs BAU — must NOT overlap caveat (5.22–5.44)
    add_text(slide, 0.4, y + 4.02, 9.2, 0.42, comparison_line, size=10, color=AMBER,
             bold=True)


def build_case_a_matched(prs):
    """Case A: Matched (Load = Gen). 24h balanced at 4,200 kWh/h."""
    layout = prs.slide_layouts[15]  # 1_Blank Slide
    slide = prs.slides.add_slide(layout)
    add_header(slide, "CANONICAL CASE A · MATCHED (Load = Gen)",
               "5-line EVN bill + CfD settlement when factory consumption exactly matches contracted solar")

    # 24h math: 4,200 kWh/h × 24h = 100,800 kWh
    load = 4200
    matched_kwh = load * HOURS  # 100,800
    # Line 1: matched * FMP * loss = market energy cost
    line1 = matched_kwh * FMP * LOSS
    # Line 2: matched * dppaCharge = service fee
    line2 = matched_kwh * FEES
    # Line 3: difference/balancing fee (assumed bundled into FEES for the teaching model;
    # deck uses 360 + 163.3 = 523.3 lumped). Show the same lumped number.
    line3 = 0  # already included in line 2
    # Line 4: shortfall retail = 0 (no shortfall)
    line4 = 0
    cevn = line1 + line2 + line3 + line4
    # Line 5: CfD = (Strike - FMP) * matched = (2000 - 1427) * 100800
    cfd = (STRIKE - FMP) * matched_kwh
    net = cevn + cfd
    # BAU: matched * retail
    bau = matched_kwh * RETAIL
    savings = bau - net

    narrative = (
        f"Factory: 4,200 kWh/h constant load × 24 h = {matched_kwh:,} kWh/day. "
        f"Solar: 4,200 kWh/h peak output × 24 h = {matched_kwh:,} kWh/day. "
        f"Every kWh of factory load is matched by a kWh of contracted solar."
    )

    evn_lines = [
        {"text": f"1. Market energy: {matched_kwh:,} kWh x {FMP:,} VND/kWh x {LOSS:.4f} = {line1:,.0f} VND", "size": 9},
        {"text": f"2. System service + DPPA fees: {matched_kwh:,} kWh x {FEES:.1f} VND/kWh = {line2:,.0f} VND", "size": 9},
        {"text": f"3. Balancing fee: bundled in line 2", "size": 9, "color": MUTED},
        {"text": f"4. Shortfall retail purchase: 0 kWh (no shortfall) = 0 VND", "size": 9, "color": MUTED},
        {"text": f"", "size": 6},
        {"text": f"Total EVN: {cevn:,.0f} VND", "size": 11, "bold": True, "color": CYAN},
    ]
    cfd_lines = [
        {"text": f"5. CfD = (Strike - FMP) x Q_cfD", "size": 9, "bold": True, "color": AMBER},
        {"text": f"     = ({STRIKE:,} - {FMP:,}) x {matched_kwh:,} kWh", "size": 9},
        {"text": f"     = +{cfd:,.0f} VND  (factory pays developer)", "size": 9, "color": AMBER},
        {"text": f"", "size": 6},
        {"text": f"FMP < Strike => CfD is a net cost to the factory at this strike level.", "size": 9, "color": MUTED},
    ]
    premium_pct = abs(savings) / bau * 100
    net_line = f"Net = EVN + CfD = {cevn:,.0f} + {cfd:,.0f} = {net:,.0f} VND  ({net/matched_kwh:,.2f} VND/kWh)"
    formula_lines = [
        {"text": f"EVN = FMP x Kpp x matched + FEES x matched = 1,427 x 1.0342 x 100,800 + 523.3 x 100,800 = {cevn:,.0f} VND", "size": 8},
        {"text": f"CfD = (Strike - FMP) x matched = (2,000 - 1,427) x 100,800 = {cfd:,.0f} VND  (positive => factory pays developer)", "size": 8},
        {"text": f"Net = {net:,.0f} VND  vs BAU = {bau:,.0f} VND  => Year-1 premium = {abs(savings):,.0f} VND (+{premium_pct:.1f}%). Savings build as EVN escalates ~4%/yr.", "size": 8, "color": MINT, "bold": True},
    ]
    comparison_line = (
        f"vs BAU retail-only: {bau:,.0f} VND ({bau/matched_kwh:,.0f} VND/kWh). "
        f"DPPA net {net:,.0f} VND ({net/matched_kwh:,.0f} VND/kWh) -- "
        f"a {premium_pct:.1f}% Year-1 premium. "
        f"The payoff is price certainty over the horizon; savings build as EVN tariffs escalate (~4%/yr)."
    )
    add_case_content(slide,
                     case_label="MATCHED  (=)",
                     header_color=MINT,
                     narrative=narrative,
                     evn_lines=evn_lines, cfd_lines=cfd_lines,
                     net_line=net_line, formula_lines=formula_lines,
                     comparison_line=comparison_line)
    add_caveat_strip(slide)
    return slide


def build_case_b_shortfall(prs):
    """Case B: Shortfall (Load > Gen). 6,100 kWh/h load vs 4,200 kWh/h gen."""
    layout = prs.slide_layouts[15]
    slide = prs.slides.add_slide(layout)
    add_header(slide, "CANONICAL CASE B · SHORTFALL (Load > Gen)",
               "5-line EVN bill + residual retail purchase + CfD when factory load exceeds contracted generation")

    load_h = 6100
    gen_h = 4200
    shortfall_h = load_h - gen_h
    matched_kwh = gen_h * HOURS  # 100,800
    shortfall_kwh = shortfall_h * HOURS  # 45,600
    total_load_kwh = load_h * HOURS  # 146,400

    line1 = matched_kwh * FMP * LOSS
    line2 = matched_kwh * FEES
    line4 = shortfall_kwh * RETAIL  # retail on shortfall
    cevn = line1 + line2 + line4
    cfd = (STRIKE - FMP) * matched_kwh  # still on matched
    net = cevn + cfd
    bau = total_load_kwh * RETAIL
    savings = bau - net

    narrative = (
        f"Factory: {load_h:,} kWh/h × 24 h = {total_load_kwh:,} kWh/day. "
        f"Solar: {gen_h:,} kWh/h × 24 h = {matched_kwh:,} kWh/day. "
        f"Shortfall: {shortfall_h:,} kWh/h × 24 h = {shortfall_kwh:,} kWh/day — bought at retail."
    )

    evn_lines = [
        {"text": f"1. Market energy: {matched_kwh:,} kWh × {FMP:,} × {LOSS:.4f} = {line1:,.0f} VND", "size": 9},
        {"text": f"2. System service + DPPA fees: {matched_kwh:,} kWh × {FEES:.1f} = {line2:,.0f} VND", "size": 9},
        {"text": f"3. Difference / balancing fee: bundled in line 2", "size": 9, "color": MUTED},
        {"text": f"4. Shortfall retail purchase: {shortfall_kwh:,} kWh × {RETAIL:,.2f} = {line4:,.0f} VND", "size": 9, "color": RED, "bold": True},
        {"text": f"", "size": 6},
        {"text": f"Total EVN: {cevn:,.0f} VND", "size": 11, "bold": True, "color": CYAN},
    ]
    cfd_lines = [
        {"text": f"5. CfD = (Strike − FMP) × Q_cfD", "size": 9, "bold": True, "color": AMBER},
        {"text": f"     = ({STRIKE:,} − {FMP:,}) × {matched_kwh:,} kWh", "size": 9},
        {"text": f"     = +{cfd:,.0f} VND  (factory pays developer)", "size": 9, "color": AMBER},
        {"text": f"", "size": 6},
        {"text": f"CfD settles only on matched kWh — shortfall is a separate retail line.", "size": 9, "color": MUTED},
        {"text": f"Two bills: DPPA + residual retail on the gap.", "size": 9, "color": MUTED},
    ]
    net_line = f"Net = EVN + CfD = {cevn:,.0f} + {cfd:,.0f} = {net:,.0f} VND  ({net/total_load_kwh:,.2f} VND/kWh blended)"
    premium_pct = abs(savings) / bau * 100
    formula_lines = [
        {"text": f"EVN = (FMP x Kpp + FEES) x matched + RETAIL x shortfall = (1,427 x 1.0342 + 523.3) x 100,800 + 2,204 x 45,600 = {cevn:,.0f} VND", "size": 8},
        {"text": f"CfD = (Strike - FMP) x matched = (2,000 - 1,427) x 100,800 = {cfd:,.0f} VND", "size": 8},
        {"text": f"Net = {net:,.0f} VND  vs BAU = {bau:,.0f} VND  => Year-1 premium = {abs(savings):,.0f} VND (+{premium_pct:.1f}%); shortfall at retail (2,204) dilutes the matched-kWh premium.", "size": 8, "color": MINT, "bold": True},
    ]
    comparison_line = (
        f"vs BAU retail-only: {bau:,.0f} VND ({bau/total_load_kwh:,.0f} VND/kWh blended). "
        f"DPPA net {net:,.0f} VND ({net/total_load_kwh:,.0f} VND/kWh blended) -- "
        f"a {premium_pct:.1f}% Year-1 premium. "
        f"Blended rate is lower than Cases A/C because shortfall kWh buy at retail (2,204). "
        f"Payoff is still price certainty; savings build as EVN escalates."
    )
    add_case_content(slide,
                     case_label="SHORTFALL  (>)",
                     header_color=AMBER,
                     narrative=narrative,
                     evn_lines=evn_lines, cfd_lines=cfd_lines,
                     net_line=net_line, formula_lines=formula_lines,
                     comparison_line=comparison_line)
    add_caveat_strip(slide)
    return slide


def build_case_c_excess(prs):
    """Case C: Excess (Load < Gen). 2,600 kWh/h load vs 4,700 kWh/h gen."""
    layout = prs.slide_layouts[15]
    slide = prs.slides.add_slide(layout)
    add_header(slide, "CANONICAL CASE C · EXCESS (Load < Gen)",
               "5-line EVN bill + CfD when contracted generation exceeds factory consumption; CfD caps at consumed volume")

    load_h = 2600
    gen_h = 4700
    excess_h = gen_h - load_h
    matched_kwh = load_h * HOURS  # 62,400 (factory is the limiting side)
    excess_kwh = excess_h * HOURS  # 50,400
    total_load_kwh = load_h * HOURS  # 62,400

    line1 = matched_kwh * FMP * LOSS
    line2 = matched_kwh * FEES
    line4 = 0  # no shortfall
    cevn = line1 + line2 + line4
    cfd = (STRIKE - FMP) * matched_kwh  # only on consumed; excess earns generator spot only
    net = cevn + cfd
    bau = total_load_kwh * RETAIL
    savings = bau - net
    excess_revenue_to_generator = excess_kwh * FMP  # generator's spot revenue on excess

    narrative = (
        f"Factory: {load_h:,} kWh/h × 24 h = {total_load_kwh:,} kWh/day. "
        f"Solar: {gen_h:,} kWh/h × 24 h = {matched_kwh + excess_kwh:,} kWh/day. "
        f"Excess: {excess_h:,} kWh/h × 24 h = {excess_kwh:,} kWh/day — generator keeps the spot revenue, no CfD applies."
    )

    evn_lines = [
        {"text": f"1. Market energy: {matched_kwh:,} kWh × {FMP:,} × {LOSS:.4f} = {line1:,.0f} VND", "size": 9},
        {"text": f"2. System service + DPPA fees: {matched_kwh:,} kWh × {FEES:.1f} = {line2:,.0f} VND", "size": 9},
        {"text": f"3. Difference / balancing fee: bundled in line 2", "size": 9, "color": MUTED},
        {"text": f"4. Shortfall retail purchase: 0 kWh (Load < Gen) = 0 VND", "size": 9, "color": MUTED},
        {"text": f"", "size": 6},
        {"text": f"Total EVN: {cevn:,.0f} VND", "size": 11, "bold": True, "color": CYAN},
    ]
    cfd_lines = [
        {"text": f"5. CfD = (Strike − FMP) × Q_cfD, capped at consumed volume", "size": 9, "bold": True, "color": AMBER},
        {"text": f"     = ({STRIKE:,} − {FMP:,}) × {matched_kwh:,} kWh", "size": 9},
        {"text": f"     = +{cfd:,.0f} VND  (factory pays developer on consumed only)", "size": 9, "color": AMBER},
        {"text": f"", "size": 6},
        {"text": f"Excess {excess_kwh:,} kWh/day → generator spot only, no CfD.", "size": 9, "color": MUTED},
        {"text": f"Generator's spot revenue on excess: ~{excess_revenue_to_generator:,.0f} VND (not your bill).", "size": 9, "color": MUTED},
    ]
    net_line = f"Net = EVN + CfD = {cevn:,.0f} + {cfd:,.0f} = {net:,.0f} VND  ({net/total_load_kwh:,.2f} VND/kWh)"
    premium_pct = abs(savings) / bau * 100
    formula_lines = [
        {"text": f"EVN = (FMP x Kpp + FEES) x matched = (1,427 x 1.0342 + 523.3) x 62,400 = {cevn:,.0f} VND", "size": 8},
        {"text": f"CfD = (Strike - FMP) x matched = (2,000 - 1,427) x 62,400 = {cfd:,.0f} VND  (capped at consumed volume)", "size": 8},
        {"text": f"Net = {net:,.0f} VND  vs BAU = {bau:,.0f} VND  => Year-1 premium = {abs(savings):,.0f} VND (+{premium_pct:.1f}%). Fees only on matched kWh; CfD capped at consumed volume.", "size": 8, "color": MINT, "bold": True},
    ]
    comparison_line = (
        f"vs BAU retail-only: {bau:,.0f} VND ({bau/total_load_kwh:,.0f} VND/kWh). "
        f"DPPA net {net:,.0f} VND ({net/total_load_kwh:,.0f} VND/kWh) -- "
        f"a {premium_pct:.1f}% Year-1 premium. "
        f"Excess generation earns developer spot only (no CfD); factory is not billed for it."
    )
    add_case_content(slide,
                     case_label="EXCESS  (<)",
                     header_color=CYAN,
                     narrative=narrative,
                     evn_lines=evn_lines, cfd_lines=cfd_lines,
                     net_line=net_line, formula_lines=formula_lines,
                     comparison_line=comparison_line)
    add_caveat_strip(slide)
    return slide


def _find_sldId_for_slide(prs, slide):
    """Return the <p:sldId> element for the given slide, or None."""
    sldIdLst = prs.slides._sldIdLst
    for child in sldIdLst:
        rid = child.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
        try:
            rel = prs.part.rels[rid]
        except KeyError:
            continue
        target = getattr(rel, "target_part", None)
        if target is slide.part:
            return child
    return None


def reorder_slide(prs, slide, target_idx_0based):
    """Move a slide to a specific 0-based position in the slide list."""
    sldIdLst = prs.slides._sldIdLst
    sldId_el = _find_sldId_for_slide(prs, slide)
    if sldId_el is None:
        raise RuntimeError("Could not find sldId for slide")

    sldIdLst.remove(sldId_el)
    children = list(sldIdLst)
    if target_idx_0based >= len(children):
        sldIdLst.append(sldId_el)
    else:
        sldIdLst.insert(children.index(children[target_idx_0based]), sldId_el)


def delete_slide_by_index(prs, idx_0based):
    """Delete the slide at the given 0-based position."""
    sldIdLst = prs.slides._sldIdLst
    children = list(sldIdLst)
    if idx_0based >= len(children):
        raise IndexError(f"Index {idx_0based} out of range (have {len(children)} slides)")
    target_sldId = children[idx_0based]
    rid = target_sldId.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
    prs.part.drop_rel(rid)
    sldIdLst.remove(target_sldId)


def main():
    prs = Presentation(SRC)
    print(f"Opening deck: {len(prs.slides)} slides")

    # 1. Add 3 new slides (will be appended at the end)
    s_matched = build_case_a_matched(prs)
    s_shortfall = build_case_b_shortfall(prs)
    s_excess = build_case_c_excess(prs)
    print(f"After adding 3 cases: {len(prs.slides)} slides")

    # 2. Reorder: move them to positions 37, 38, 39 (0-based) so they become slides 38, 39, 40 (1-based)
    # The workshop scenarios start at original position 37 (0-based 36); we want the new cases
    # to come right after slide 37 (1-based) = position 37 (0-based).
    reorder_slide(prs, s_matched, 37)   # becomes slide 38
    reorder_slide(prs, s_shortfall, 38)  # becomes slide 39
    reorder_slide(prs, s_excess, 39)     # becomes slide 40

    # Verify the layout by printing slide titles
    print("After reordering, slides 36-45 are:")
    for i in range(35, 45):
        slide = prs.slides[i]
        title = ''
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text.strip():
                title = shape.text.strip().split('\n')[0][:60]
                break
        print(f"  {i+1}. {title}")
    print()

    # 3. Delete the old Scenarios 1+3 slides.
    #    After reordering, layout is:
    #      0-based 0-36: original slides 1-37 (unchanged)
    #      0-based 37-39: new cases A, B, C
    #      0-based 40-43: old Scenarios 1 (originally 1-based 38-41)
    #      0-based 44-47: old Scenarios 3 (originally 1-based 42-45)
    #    So delete 0-based 40-47 (8 slides, in reverse order).
    for idx in reversed(range(40, 48)):
        delete_slide_by_index(prs, idx)
    print(f"After deleting 8 old scenario slides: {len(prs.slides)} slides (was 57+3-8=52)")

    prs.save(OUT)
    print(f"Saved to {OUT}")


if __name__ == "__main__":
    main()
