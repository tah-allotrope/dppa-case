from copy import deepcopy
from pathlib import Path

import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
REF = ROOT / "ref" / "DPPA 2025 ref.pptx"
OUT = ROOT / "dppa-2026-factory-energy-proposal.pptx"
ASSET_DIR = ROOT / "deck-qa" / "generated-2026"
ASSET_DIR.mkdir(parents=True, exist_ok=True)

COLORS = {
    "dark": "2E3A49",
    "green": "1B786E",
    "green_dark": "2F645B",
    "grey": "F7F7F7",
    "text": "131314",
    "muted": "5B646E",
    "cyan": "168A9A",
    "amber": "F1A722",
    "magenta": "8B4F96",
    "red": "C54E4B",
    "line": "D7E2E0",
}

F_HEAD = "Montserrat"
F_BODY = "Roboto"

HOURS = list(range(24))
FMP_SHAPE = [
    0.70, 0.69, 0.68, 0.69, 0.72,
    0.78, 0.84, 0.88, 0.92, 0.97,
    1.00, 1.04, 1.08, 1.11, 1.15,
    1.18, 1.22, 1.28, 1.36, 1.42,
    1.30, 1.08, 0.92, 0.80,
]

INPUTS = {
    "strike": 2100,
    "market": 1700,
    "dppa": 523.34,
    "loss": 1.027263,
    "retail": 2100,
}
FMP = [round(INPUTS["market"] * m) for m in FMP_SHAPE]


def rgb(hex_color):
    hex_color = hex_color.strip("#")
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))


def solar_curve(scale, shoulder):
    import math

    values = []
    for hour in HOURS:
        if hour < 6 or hour > 18:
            values.append(0)
        else:
            normalized = math.sin(((hour - 6) / 12) * math.pi)
            values.append(round((max(0, normalized) ** (1 + shoulder)) * scale))
    return values


SCENARIOS = {
    "higherLoad": {
        "label": "Load > Gen",
        "load": [4300 if h < 6 else 5200 if h < 10 else 6100 if h < 17 else 5000 if h < 22 else 4400 for h in HOURS],
        "gen": solar_curve(4200, 0.42),
    },
    "balanced": {
        "label": "Load = Gen",
        "load": [3000 if h < 6 else 4000 if h < 9 else 4700 if h < 16 else 3900 if h < 20 else 3200 for h in HOURS],
        "gen": solar_curve(4700, 0.28),
    },
    "higherGen": {
        "label": "Load < Gen",
        "load": [2600 if h < 7 else 3200 if h < 10 else 3600 if h < 16 else 3100 if h < 20 else 2700 for h in HOURS],
        "gen": solar_curve(6200, 0.22),
    },
}


def calc(sc, mode="matched"):
    intervals = []
    for hour in HOURS:
        load = sc["load"][hour]
        gen = sc["gen"][hour]
        matched = min(load, gen)
        shortfall = max(load - gen, 0)
        excess = max(gen - load, 0)
        contract = gen if mode == "generation" else matched
        evn_market = matched * FMP[hour] * INPUTS["loss"]
        evn_dppa = matched * INPUTS["dppa"]
        evn_retail = shortfall * INPUTS["retail"]
        evn_total = evn_market + evn_dppa + evn_retail
        developer = contract * (INPUTS["strike"] - FMP[hour])
        total = evn_total + developer
        baseline = load * INPUTS["retail"]
        intervals.append({
            "hour": hour,
            "load": load,
            "gen": gen,
            "matched": matched,
            "shortfall": shortfall,
            "excess": excess,
            "contract": contract,
            "fmp": FMP[hour],
            "evn_market": evn_market,
            "evn_dppa": evn_dppa,
            "evn_retail": evn_retail,
            "evn_total": evn_total,
            "developer": developer,
            "total": total,
            "baseline": baseline,
        })
    totals = {k: sum(iv[k] for iv in intervals) for k in ["load", "gen", "matched", "shortfall", "excess", "total", "baseline"]}
    totals["savings"] = totals["baseline"] - totals["total"]
    totals["blended"] = totals["total"] / totals["load"]
    totals["match_rate"] = totals["matched"] / totals["load"]
    return {"intervals": intervals, "totals": totals}


RESULTS = {key: calc(value) for key, value in SCENARIOS.items()}


def fmt(n):
    return f"{round(n):,}"


def clear_text(shape):
    tf = shape.text_frame
    tf.clear()
    return tf


def suppress_bullet(p):
    pPr = p._p.get_or_add_pPr()
    for tag in ("a:buChar", "a:buAutoNum", "a:buNone"):
        for el in pPr.findall(qn(tag)):
            pPr.remove(el)
    pPr.append(pPr.makeelement(qn("a:buNone"), {}))


def set_text(shape, text, size=14, bold=False, color="131314", font=F_BODY, align=None):
    tf = clear_text(shape)
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    suppress_bullet(p)
    if align is not None:
        p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = rgb(color)


def add_textbox(slide, x, y, w, h, text, size=10, bold=False, color="131314", font=F_BODY, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    set_text(box, text, size=size, bold=bold, color=color, font=font, align=align)
    return box


def add_bullets(slide, x, y, w, h, items, size=9.5, color="131314"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = F_BODY
        p.font.size = Pt(size)
        p.font.color.rgb = rgb(color)
        p.space_after = Pt(5)
    return box


def add_stat_card(slide, x, y, w, h, value, label, accent="1B786E"):
    card = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid()
    card.fill.fore_color.rgb = rgb("F7F7F7")
    card.line.color.rgb = rgb("D7E2E0")
    add_textbox(slide, x + 0.08, y + 0.12, w - 0.16, 0.27, value, 18, True, accent, F_HEAD, PP_ALIGN.CENTER)
    add_textbox(slide, x + 0.08, y + 0.48, w - 0.16, 0.24, label, 7.2, False, COLORS["muted"], F_BODY, PP_ALIGN.CENTER)
    return card


def delete_shape(shape):
    el = shape._element
    el.getparent().remove(el)


def add_picture_cover(slide, path, x, y, w, h):
    slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))


def make_chart(path, scenario, title=None, compact=False):
    plt.rcParams["font.family"] = "DejaVu Sans"
    fig_w, fig_h = (9.4, 3.5) if not compact else (6.4, 1.9)
    fig, ax = plt.subplots(figsize=(fig_w, fig_h), dpi=220)
    fig.patch.set_facecolor("white")
    ax.set_facecolor("white")
    load = scenario["load"]
    gen = scenario["gen"]
    matched = [min(a, b) for a, b in zip(load, gen)]
    ax.fill_between(HOURS, 0, matched, color=f"#{COLORS['green']}", alpha=0.16, linewidth=0)
    ax.plot(HOURS, load, color=f"#{COLORS['cyan']}", linewidth=2.4, label="Factory load")
    ax.plot(HOURS, gen, color=f"#{COLORS['amber']}", linewidth=2.4, label="Solar generation")
    ax2 = ax.twinx()
    ax2.plot(HOURS, FMP, color=f"#{COLORS['magenta']}", linewidth=1.8, label="FMP")
    ax.set_xlim(0, 23)
    ax.set_ylim(0, 6800)
    ax2.set_ylim(900, 2600)
    ax.set_xticks(range(0, 24, 3 if not compact else 6))
    ax.grid(axis="y", color="#D7E2E0", linewidth=0.7)
    ax.tick_params(axis="both", labelsize=8 if not compact else 6.5, colors="#5B646E")
    ax2.tick_params(axis="y", labelsize=8 if not compact else 6.5, colors="#5B646E")
    ax.set_ylabel("kWh / hour", fontsize=8 if not compact else 6.5, color="#5B646E")
    ax2.set_ylabel("VND/kWh", fontsize=8 if not compact else 6.5, color="#5B646E")
    if title:
        ax.set_title(title, loc="left", fontsize=11, fontweight="bold", color=f"#{COLORS['dark']}", pad=8)
    lines = ax.get_lines() + ax2.get_lines()
    labels = [line.get_label() for line in lines]
    ax.legend(lines, labels, loc="upper left", frameon=False, fontsize=8 if not compact else 6.5, ncol=3)
    for spine in ax.spines.values():
        spine.set_color("#D7E2E0")
    for spine in ax2.spines.values():
        spine.set_color("#D7E2E0")
    fig.tight_layout(pad=0.8)
    fig.savefig(path, transparent=False, bbox_inches="tight")
    plt.close(fig)


def make_scenario_table(path):
    rows = [
        ["Metric", "Load > Gen", "Load = Gen", "Load < Gen"],
        ["Daily load kWh", fmt(RESULTS["higherLoad"]["totals"]["load"]), fmt(RESULTS["balanced"]["totals"]["load"]), fmt(RESULTS["higherGen"]["totals"]["load"])],
        ["Matched kWh", fmt(RESULTS["higherLoad"]["totals"]["matched"]), fmt(RESULTS["balanced"]["totals"]["matched"]), fmt(RESULTS["higherGen"]["totals"]["matched"])],
        ["Shortfall kWh", fmt(RESULTS["higherLoad"]["totals"]["shortfall"]), fmt(RESULTS["balanced"]["totals"]["shortfall"]), fmt(RESULTS["higherGen"]["totals"]["shortfall"])],
        ["Excess kWh", fmt(RESULTS["higherLoad"]["totals"]["excess"]), fmt(RESULTS["balanced"]["totals"]["excess"]), fmt(RESULTS["higherGen"]["totals"]["excess"])],
        ["DPPA cost VND", fmt(RESULTS["higherLoad"]["totals"]["total"]), fmt(RESULTS["balanced"]["totals"]["total"]), fmt(RESULTS["higherGen"]["totals"]["total"])],
        ["BAU cost VND", fmt(RESULTS["higherLoad"]["totals"]["baseline"]), fmt(RESULTS["balanced"]["totals"]["baseline"]), fmt(RESULTS["higherGen"]["totals"]["baseline"])],
        ["Savings vs BAU", fmt(RESULTS["higherLoad"]["totals"]["savings"]), fmt(RESULTS["balanced"]["totals"]["savings"]), fmt(RESULTS["higherGen"]["totals"]["savings"])],
    ]
    fig, ax = plt.subplots(figsize=(8.8, 3.25), dpi=220)
    fig.patch.set_facecolor("white")
    ax.axis("off")
    table = ax.table(cellText=rows, cellLoc="center", loc="center", colWidths=[0.30, 0.23, 0.23, 0.24])
    table.auto_set_font_size(False)
    table.set_fontsize(8.2)
    for (r, c), cell in table.get_celld().items():
        cell.set_edgecolor("#D7E2E0")
        cell.set_linewidth(0.7)
        cell.set_height(0.105)
        if r == 0:
            cell.set_facecolor("#EAF3F1")
            cell.get_text().set_fontweight("bold")
            cell.get_text().set_color(f"#{COLORS['dark']}")
        elif c == 0:
            cell.get_text().set_ha("left")
            cell.get_text().set_fontweight("bold")
        if r == 7 and c > 0:
            cell.get_text().set_color(f"#{COLORS['red']}")
            cell.get_text().set_fontweight("bold")
    fig.tight_layout(pad=0.3)
    fig.savefig(path, transparent=False, bbox_inches="tight")
    plt.close(fig)


def make_equation_card(path):
    fig, ax = plt.subplots(figsize=(7.5, 1.55), dpi=240)
    fig.patch.set_facecolor("white")
    ax.axis("off")
    ax.text(
        0.02,
        0.72,
        "EVN Market + DPPA Charge + CfD",
        fontsize=14,
        fontweight="bold",
        color=f"#{COLORS['dark']}",
        transform=ax.transAxes,
    )
    ax.text(
        0.02,
        0.42,
        "= Q × FMP + Q × C_DPPA + Q × (Strike − FMP)",
        fontsize=12.5,
        color=f"#{COLORS['green_dark']}",
        transform=ax.transAxes,
    )
    ax.text(
        0.02,
        0.14,
        "= Q × Strike + Q × C_DPPA   (+ loss adjustment)",
        fontsize=12.5,
        color=f"#{COLORS['green_dark']}",
        transform=ax.transAxes,
    )
    ax.add_patch(plt.Rectangle((0, 0), 1, 1, fill=False, edgecolor=f"#{COLORS['green']}", linewidth=1.2, transform=ax.transAxes))
    fig.tight_layout(pad=0.2)
    fig.savefig(path, transparent=False, bbox_inches="tight")
    plt.close(fig)


def generate_assets():
    make_chart(ASSET_DIR / "balanced_chart.png", SCENARIOS["balanced"], "Scenario B: balanced daily profile")
    make_chart(ASSET_DIR / "load_gt_gen_chart.png", SCENARIOS["higherLoad"], "Scenario A: load above generation", compact=True)
    make_chart(ASSET_DIR / "gen_gt_load_chart.png", SCENARIOS["higherGen"], "Scenario C: generation above load", compact=True)
    make_scenario_table(ASSET_DIR / "scenario_table.png")
    make_equation_card(ASSET_DIR / "cancellation_card.png")


def cover(slide, x, y, w, h, color="FFFFFF"):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color)
    shape.line.color.rgb = rgb(color)
    return shape


def main():
    generate_assets()
    prs = Presentation(str(REF))

    # Slide 1
    set_text(prs.slides[0].shapes[5], "Session 4.4: Off-Site Solutions:\nAddressing DPPA Pricing Concerns", 21, True, "F7F7F7", F_HEAD)

    # Slide 2: replace speaker content with update overview, preserve true logos/photo frame.
    s = prs.slides[1]
    set_text(s.shapes[3], "About the 2026 Update", 25, True, COLORS["dark"], F_HEAD)
    set_text(
        s.shapes[5],
        "Current repo: CFO calculator, tested settlement engine, formula docs, and workshop readiness analysis",
        12.5,
        True,
        COLORS["green"],
        F_HEAD,
    )
    set_text(
        s.shapes[4],
        "The reference deck explained why DPPA buyers worry about market price. This update uses the current repository to show the actual settlement mechanics behind that concern.\n\n"
        "Updated storyline:\n"
        "• FMP and CFMP remain visible buyer concerns\n"
        "• Hourly matched quantity is the cost driver\n"
        "• The CfD cancels much of the FMP exposure on matched kWh\n"
        "• Settlement quantity mode determines overgeneration risk\n\n"
        "Current app status: 35 tests passing, three synthetic profiles, default matched-mode settlement.",
        9.6,
        False,
        COLORS["text"],
        F_BODY,
    )
    # Remove speaker portrait; replace with a clean KPI panel instead of a dark screenshot.
    delete_shape(s.shapes[7])
    cover(s, 1.02, 1.0, 3.75, 2.55, "FFFFFF")
    add_stat_card(s, 1.28, 1.28, 1.35, 0.86, "35", "tests passing")
    add_stat_card(s, 3.0, 1.28, 1.35, 0.86, "3", "factory profiles")
    add_stat_card(s, 1.28, 2.45, 1.35, 0.86, "5", "pricing inputs")
    add_stat_card(s, 3.0, 2.45, 1.35, 0.86, "0", "critical gaps")

    # Slide 3
    s = prs.slides[2]
    set_text(s.shapes[4], "Understanding Market Price", 22, True, COLORS["dark"], F_HEAD)
    set_text(
        s.shapes[3],
        "Spot market price remains the visible buyer concern, but the current calculator shows where it truly affects cost",
        18,
        True,
        COLORS["text"],
        F_HEAD,
    )
    delete_shape(s.shapes[2])
    add_picture_cover(s, ASSET_DIR / "balanced_chart.png", 1.06, 1.92, 8.48, 3.15)

    # Slide 4
    s = prs.slides[3]
    set_text(s.shapes[4], "Understanding Market Price", 22, True, COLORS["dark"], F_HEAD)
    set_text(
        s.shapes[2],
        "Matched quantity, not headline FMP, drives the buyer's DPPA cost exposure",
        15,
        True,
        COLORS["text"],
        F_HEAD,
    )
    s.shapes[2].height = Inches(0.68)
    delete_shape(s.shapes[3])
    cover(s, 1.14, 1.42, 6.9, 3.75)
    stats = [
        ("2,100", "strike price Pc"),
        ("523.34", "DPPA charge"),
        ("1.027263", "Kpp loss factor"),
    ]
    for i, (value, label) in enumerate(stats):
        x = 1.35 + i * 2.2
        add_textbox(s, x, 2.05, 1.75, 0.35, value, 20, True, COLORS["green"], F_HEAD, PP_ALIGN.CENTER)
        add_textbox(s, x, 2.45, 1.75, 0.22, label, 8.5, False, COLORS["muted"], F_BODY, PP_ALIGN.CENTER)
    add_bullets(
        s,
        1.35,
        3.18,
        6.0,
        1.15,
        [
            "FMP remains visible in the EVN market charge and developer CfD.",
            "Matched kWh mostly collapse toward strike price plus DPPA charge.",
            "Shortfall and excess volumes decide whether that clean story holds.",
        ],
        10,
    )

    # Slide 5
    s = prs.slides[4]
    set_text(s.shapes[0], "Understanding Market Price", 22, True, COLORS["dark"], F_HEAD)
    set_text(
        s.shapes[3],
        "The repo uses a synthetic daily FMP curve to teach price sensitivity without hiding the settlement logic.",
        17.5,
        True,
        COLORS["text"],
        F_HEAD,
    )
    delete_shape(s.shapes[5])
    delete_shape(s.shapes[4])
    add_picture_cover(s, ASSET_DIR / "load_gt_gen_chart.png", 1.05, 2.14, 4.47, 2.55)
    add_picture_cover(s, ASSET_DIR / "gen_gt_load_chart.png", 5.72, 2.14, 4.05, 2.55)

    # Slide 6 payment mechanisms: keep reference diagram structure, update formula language.
    s = prs.slides[5]
    set_text(s.shapes[1], "Payment Mechanisms", 22, True, COLORS["dark"], F_HEAD)
    set_text(s.shapes[9], "Buyer cost = EVN market + DPPA charge + retail shortfall + developer CfD", 13.5, True, COLORS["green_dark"], F_HEAD)
    set_text(s.shapes[11], "Developer CfD\n= Contract quantity × (Strike − FMP)", 9.2, False, COLORS["green_dark"], F_BODY)
    set_text(s.shapes[12], "DPPA service charge\n= Matched kWh × C_DPPA", 9.2, False, COLORS["green_dark"], F_BODY)
    set_text(s.shapes[13], "Retail shortfall\n= (Load − matched kWh) × retail tariff", 9.2, False, COLORS["green_dark"], F_BODY)
    set_text(s.shapes[14], "EVN market\n= Matched kWh × FMP × Kpp", 9.2, False, COLORS["green_dark"], F_BODY)
    set_text(s.shapes[6], "FMP reference\ninside EVN bill", 9.2, True, COLORS["green_dark"], F_HEAD)
    set_text(s.shapes[15], "Strike price − spot market", 10.5, True, COLORS["text"], F_HEAD)

    # Slide 7 chart-focused diurnal profile.
    s = prs.slides[6]
    set_text(s.shapes[4], "Diurnal Profile", 22, True, COLORS["dark"], F_HEAD)
    set_text(s.shapes[2], "Accurate allocation of QKHhc is critical for determining DPPA cost settlement", 15, True, COLORS["text"], F_HEAD)
    delete_shape(s.shapes[3])
    add_picture_cover(s, ASSET_DIR / "balanced_chart.png", 0.99, 1.74, 8.78, 3.55)

    # Slide 8: formulas/definitions.
    s = prs.slides[7]
    set_text(s.shapes[0], "QKHhc is the key factor in all DPPA cost calculations: buyers must calculate it accurately.", 21, True, COLORS["dark"], F_HEAD)
    set_text(s.shapes[3], "", 9.2, False, COLORS["text"], F_BODY)
    delete_shape(s.shapes[7])
    delete_shape(s.shapes[6])
    formulas = [
        ("Matched volume", "QKHhc(i) = min(QKH(i), Qm(i))"),
        ("Shortfall", "max(QKH(i) − Qm(i), 0)"),
        ("Excess", "max(Qm(i) − QKH(i), 0)"),
    ]
    for i, (head, body) in enumerate(formulas):
        x = 1.05 + i * 2.75
        cover(s, x, 1.35, 2.35, 0.92, "F7F7F7")
        add_textbox(s, x + 0.13, 1.52, 2.1, 0.18, head, 9.5, True, COLORS["green_dark"], F_HEAD)
        add_textbox(s, x + 0.13, 1.82, 2.1, 0.22, body, 8.8, False, COLORS["text"], F_BODY)
    add_bullets(
        s,
        1.05,
        2.78,
        7.9,
        1.2,
        [
            "Matched mode: contract quantity equals matched energy.",
            "Generation mode: contract quantity equals renewable generation, so excess solar can still create CfD exposure.",
            "Allocated mode: contract quantity follows the allocation in the contract.",
        ],
        10,
    )
    add_textbox(
        s,
        1.05,
        4.45,
        7.8,
        0.36,
        "Buyer implication: negotiate the quantity definition before debating a headline strike price.",
        11,
        True,
        COLORS["green_dark"],
        F_HEAD,
    )

    # Slide 9: CFD
    s = prs.slides[8]
    set_text(s.shapes[5], "Payment Mechanism under the CFD", 22, True, COLORS["dark"], F_HEAD)
    set_text(
        s.shapes[4],
        "Strike price determines financial settlement through CFD, converting market price volatility into predictable financial flows",
        17.5,
        True,
        COLORS["text"],
        F_HEAD,
    )
    set_text(s.shapes[2], "", 9.2, False, COLORS["text"], F_BODY)
    delete_shape(s.shapes[3])
    add_picture_cover(s, ASSET_DIR / "cancellation_card.png", 1.25, 1.88, 6.8, 1.45)
    add_bullets(
        s,
        1.28,
        3.7,
        7.4,
        0.82,
        [
            "Pc(i): committed strike price at cycle i.",
            "FMP(i): hourly market reference price.",
            "Qca(i): contract quantity; default app mode equals matched volume.",
        ],
        9.2,
    )

    # Slide 10: feasibility table
    s = prs.slides[9]
    set_text(s.shapes[4], "Quantify DPPA Costs to Assess Feasibility & Negotiate RE Supply", 20, True, COLORS["dark"], F_HEAD)
    set_text(s.shapes[3], "Feasibility analysis enables buyers to compare DPPA against traditional procurement with confidence", 14, True, COLORS["text"], F_HEAD)
    delete_shape(s.shapes[0])
    add_picture_cover(s, ASSET_DIR / "scenario_table.png", 1.28, 1.78, 8.39, 3.35)

    # Slide 11 section transition.
    set_text(prs.slides[10].shapes[5], "Session 4.5: Preparing for DPPA Implementation", 23, True, "F7F7F7", F_HEAD)

    # Slide 12 exercise.
    s = prs.slides[11]
    set_text(s.shapes[5], "Interactive Exercise:\nDPPA Scenario Analysis  (~60 minutes)", 23, True, COLORS["grey"], F_HEAD)
    set_text(
        s.shapes[6],
        "Use the CFO calculator to test one factory load profile against three questions:\n\n"
        "• How much load is matched by renewable generation?\n"
        "• Which hours create retail shortfall exposure?\n"
        "• Which settlement quantity mode is acceptable before negotiating strike price?\n\n"
        "Start with the balanced profile, then compare Load > Gen and Load < Gen.",
        12,
        False,
        COLORS["grey"],
        F_BODY,
    )

    # Slide 13 negotiation prompt.
    s = prs.slides[12]
    set_text(s.shapes[9], "Scenario DPPA Negotiation", 22, True, COLORS["dark"], F_HEAD)
    set_text(s.shapes[2], "Scenario DPPA Negotiation", 18, True, COLORS["dark"], F_HEAD)
    set_text(
        s.shapes[4],
        "Your factory is located in Northern Vietnam, within an industrial park, and is currently purchasing electricity from an EVN-subsidiary or industrial park electricity retailer. Annual demand is 150 GWh.\n\n"
        "Determine the essential conditions for proceeding with an appropriate DPPA arrangement:\n\n"
        "• Electricity retailer\n"
        "• Optimal renewable energy source\n"
        "• Required capacity and load match\n"
        "• DPPA contract duration\n"
        "• Settlement quantity mode\n"
        "• Risk allocation for shortfall and excess",
        9.4,
        False,
        COLORS["text"],
        F_BODY,
    )
    delete_shape(s.shapes[8])
    delete_shape(s.shapes[7])
    add_picture_cover(s, ASSET_DIR / "balanced_chart.png", 3.86, 1.17, 6.04, 2.88)

    prs.save(str(OUT))
    print(f"Saved {OUT}")


if __name__ == "__main__":
    main()
